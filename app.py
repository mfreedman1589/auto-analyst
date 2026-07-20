import streamlit as st
import streamlit.components.v1 as components
import pandas as pd
import requests
from requests.adapters import HTTPAdapter
from urllib3.util.retry import Retry
import re
import concurrent.futures
import plotly.express as px
from urllib.parse import urlparse
from bs4 import BeautifulSoup
import datetime
import pytz
from fpdf import FPDF
import io
import matplotlib
matplotlib.use("Agg")   # headless rendering for chart images
import matplotlib.pyplot as plt
import os
import html
import random
import threading
import time
from streamlit_gsheets import GSheetsConnection

# --- CONFIGURATION & CUSTOM CSS ---
st.set_page_config(page_title="Auto-Sales Intelligence Agent", layout="wide")

st.markdown("""
    <style>
        [data-testid="stSidebarHeader"] { position: sticky !important; top: 0 !important; z-index: 999 !important; background-color: var(--secondary-background-color) !important; }
        [data-testid="stSidebarResizer"] { height: 100vh !important; }
        button[kind="primary"] { background-color: #D70015 !important; color: white !important; border: 1px solid #A30010 !important; font-weight: bold !important; border-radius: 6px !important; }
        button[kind="primary"]:hover { background-color: #FF3B30 !important; border: 1px solid #D70015 !important; color: white !important; }
    </style>
""", unsafe_allow_html=True)

# --- SESSION STATE INITIALIZATION ---
if 'history' not in st.session_state:
    st.session_state.history = {} 
if 'current_report_id' not in st.session_state:
    st.session_state.current_report_id = None 
if 'min_visitors' not in st.session_state:
    st.session_state.min_visitors = 1
if 'global_usage_count' not in st.session_state:
    st.session_state.global_usage_count = "..."
if 'mc_month_usage' not in st.session_state:
    st.session_state.mc_month_usage = None

# --- THE DEALER API VAULT (Google Sheets Powered) ---
FALLBACK_VAULT = {
    'hananiavw.com': {'app_id': 'YL5AFXM3DW', 'api_key': '59d32b7b5842f84284e044c7ca465498', 'index': 'volkswagenoforangepark-sbm0424_production_inventory'},
    'hondasanmarcos.com': {'app_id': 'V3ZOVI2QFZ', 'api_key': 'ec7553dd56e6d4c8bb447a0240e7aab3', 'index': 'hondaofsanmarcos_production_inventory'}
}

def _clean_cell(value):
    """Normalize a sheet cell to a stripped string ('' for blanks/NaN)."""
    text = str(value).strip()
    return "" if text.lower() in ("", "nan", "none") else text

def _row_is_ignored(app_id, api_key, index_name):
    """A domain is IGNORE if any credential column is literally IGNORE (case-insensitive)."""
    return any(v.upper() == "IGNORE" for v in (app_id, api_key, index_name))

def build_vault_dict(vault_df):
    """
    Turn the sheet DataFrame into the in-memory vault.
    Ignored domains are stored as {'ignore': True} so the scanner has a single,
    foolproof flag to check. Domains with missing creds are simply left out
    (they fall through to the visual scraper).
    """
    vault_dict = {}
    for _, row in vault_df.iterrows():
        domain = _clean_cell(row.get('Base Domain', ''))
        if not domain:
            continue
        app_id = _clean_cell(row.get('App ID', ''))
        api_key = _clean_cell(row.get('API Key', ''))
        index_name = _clean_cell(row.get('Index Name', ''))

        if _row_is_ignored(app_id, api_key, index_name):
            vault_dict[domain] = {'ignore': True}          # hard kill switch
        elif app_id and api_key and index_name:
            vault_dict[domain] = {
                'app_id': app_id,
                'api_key': api_key,
                'index': index_name,
            }
        # else: incomplete row -> not added -> scraper handles it
    return vault_dict

def load_vault():
    """
    The Google Sheet is the ONLY source of truth. Returns:
      (vault_dict, vault_df, conn, loaded_ok)
    On failure we fall back to FALLBACK_VAULT but flag loaded_ok=False so the UI
    can warn the user rather than silently running on stale/tiny credentials.
    """
    try:
        conn = st.connection("gsheets", type=GSheetsConnection)
        vault_df = conn.read(worksheet="Sheet1", ttl=60)
        vault_dict = build_vault_dict(vault_df)

        try:
            usage_df = conn.read(worksheet="UsageStats", ttl=60)
            st.session_state.global_usage_count = len(usage_df)
            # Sum this month's market lookups from our own report log, so the
            # sidebar can show remaining budget without spending an API call.
            if "MarketCheck Calls" in usage_df.columns:
                ts = pd.to_datetime(usage_df["Timestamp"], errors="coerce", utc=True)
                now_ts = pd.Timestamp.now(tz="UTC")
                in_month = (ts.dt.year == now_ts.year) & (ts.dt.month == now_ts.month)
                st.session_state.mc_month_usage = int(
                    pd.to_numeric(usage_df.loc[in_month, "MarketCheck Calls"],
                                  errors="coerce").fillna(0).sum()
                )
        except Exception:
            pass

        return vault_dict, vault_df, conn, True
    except Exception:
        return dict(FALLBACK_VAULT), None, None, False

DEALER_API_VAULT, vault_df, gsheets_conn, VAULT_LOADED_OK = load_vault()

if not VAULT_LOADED_OK:
    st.warning(
        "Couldn't load saved dealer settings right now — some dealers may not "
        "resolve, and new dealers can't be saved until the connection is back.",
        icon="⚠️",
    )

def upsert_vault_row(domain, app_id, api_key, index_name):
    """
    Single write path for the vault. Keeps the sheet authoritative:
      1. Re-read the live sheet (ttl=0, no cache)
      2. Drop any existing row(s) for this domain (prevents duplicate/ghost rows)
      3. Append the new row
      4. Write back and bust the cache so the next load re-reads fresh
    Returns True on success. There is NO local session overlay — after this
    returns, the app reruns and reloads straight from the sheet.
    """
    if gsheets_conn is None:
        return False
    try:
        live_df = gsheets_conn.read(worksheet="Sheet1", ttl=0)
        live_df = live_df[live_df['Base Domain'].astype(str).str.strip() != '']
        # Remove existing entries for this domain (case-insensitive) to avoid dupes
        keep_mask = live_df['Base Domain'].astype(str).str.strip().str.lower() != domain.lower()
        live_df = live_df[keep_mask]
        new_row = pd.DataFrame([{
            'Base Domain': domain,
            'App ID': app_id,
            'API Key': api_key,
            'Index Name': index_name,
        }])
        updated_df = pd.concat([live_df, new_row], ignore_index=True)
        gsheets_conn.update(worksheet="Sheet1", data=updated_df)
        st.cache_data.clear()
        return True
    except Exception:
        return False

# --- VAULT ONBOARDING HELPERS (salesperson-friendly) ---
def normalize_domain(raw):
    """Strip protocol / www / paths down to a bare base domain."""
    d = str(raw).strip().lower()
    d = re.sub(r'^https?://', '', d)
    d = d.replace('www.', '')
    d = d.split('/')[0].split('?')[0].strip().strip('/')
    return d

def parse_algolia_credentials(text, fallback_domain=""):
    """
    Forgiving extractor. Accepts ANY of:
      • Bookmarklet output:   domain|APPID|APIKEY|indexname
      • A raw Algolia URL:    ...x-algolia-application-id=APPID&x-algolia-api-key=APIKEY.../indexes/NAME/query
      • Messy pasted config:  loose text containing the three values
    Returns dict(domain, app_id, api_key, index) or None if it can't find all three.
    """
    if not text:
        return None
    text = str(text).strip()
    domain, app_id, api_key, index = fallback_domain, "", "", ""

    # 1) Pipe format from the bookmarklet: domain|appId|apiKey|index
    if text.count('|') >= 3:
        parts = [p.strip() for p in text.split('|')]
        domain = parts[0] or domain
        app_id, api_key, index = parts[1], parts[2], parts[3]

    # 2) URL query-param format (the old DevTools method)
    if not app_id:
        m = re.search(r'x-algolia-application-id=([A-Za-z0-9]+)', text)
        if m: app_id = m.group(1)
    if not api_key:
        m = re.search(r'x-algolia-api-key=([^&\s"\']{20,})', text, re.I)
        if m: api_key = m.group(1)
    if not index:
        m = re.search(r'/indexes/([^/?"\']+)/(?:query|queries)', text)
        if m: index = m.group(1)

    # 3) Loose fallbacks — pull the values out of raw page text / config
    if not app_id:
        m = (re.search(r'([A-Z0-9]{10})-dsn\.algolia\.net', text)
             or re.search(r'(?:application-?id|appid|app_id)["\'\s:=]+([A-Z0-9]{10})', text, re.I))
        if m: app_id = m.group(1).upper()
    if not api_key:
        m = re.search(r'(?:x-algolia-api-key|algolia[_-]?api[_-]?key|api[_-]?key|apikey|searchapikey|search[_-]?key)["\'\s:=]+["\']?([A-Za-z0-9+/=_-]{20,})', text, re.I)
        if m: api_key = m.group(1)
    if not index:
        m = (re.search(r'["\']?index(?:name)?["\']?\s*[:=]\s*["\']([a-z0-9_\-]*inventory[a-z0-9_\-]*)["\']', text, re.I)
             or re.search(r'([a-z0-9][a-z0-9_\-]*_production_inventory[a-z0-9_]*)', text, re.I))
        if m: index = m.group(1)

    if app_id and api_key and index:
        return {
            'domain': normalize_domain(domain) if domain else "",
            'app_id': app_id.strip(),
            'api_key': api_key.strip(),
            'index': index.strip(),
        }
    return None

def autodetect_algolia(domain):
    """
    Server-side attempt: fetch the dealer's pages and scrape the Algolia creds
    out of the page config. An aggressive firewall can block this — in which case
    the salesperson should use the Bookmarklet (runs in their browser).
    Returns a creds dict or None.
    """
    domain = normalize_domain(domain)
    if not domain:
        return None
    candidate_paths = ['', '/new-inventory/', '/used-inventory/', '/inventory/', '/searchnew.aspx']
    blob = ""
    for path in candidate_paths:
        url = f"https://www.{domain}{path}"
        try:
            r = requests.get(url, timeout=8, allow_redirects=True,
                             headers={'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) '
                                                    'AppleWebKit/537.36 (KHTML, like Gecko) '
                                                    'Chrome/122.0.0.0 Safari/537.36'})
            if r.status_code == 200 and r.text:
                blob += "\n" + r.text
                creds = parse_algolia_credentials(blob, fallback_domain=domain)
                if creds:  # found everything — stop early
                    return creds
        except Exception:
            continue
    return parse_algolia_credentials(blob, fallback_domain=domain)

def test_algolia_creds(app_id, api_key, index, domain=""):
    """
    Fire ONE Algolia request against the given creds to diagnose them.
    Sends the same Referer/Origin as scan_url so a referer-locked key is tested
    the way it will actually be used. Returns (ok, headline, detail) in plain
    English so a salesperson can see whether the keys work.
    """
    if not (app_id and api_key and index):
        return False, "Missing keys", "One of App ID / API Key / Index is blank."
    endpoint = f"https://{app_id.lower()}-dsn.algolia.net/1/indexes/{index}/query"
    headers = {
        "x-algolia-application-id": app_id,
        "x-algolia-api-key": api_key,
        "Content-Type": "application/json",
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                      "(KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36",
    }
    if domain:
        d = normalize_domain(domain)
        headers["Referer"] = f"https://www.{d}/"
        headers["Origin"] = f"https://www.{d}"
    try:
        r = requests.post(endpoint, headers=headers,
                          json={"params": "query=&hitsPerPage=1"}, timeout=8)
    except Exception as e:
        return False, "Couldn't reach Algolia", str(e)[:200]

    if r.status_code == 200:
        try:
            data = r.json()
            n = data.get("nbHits", 0)
        except Exception:
            return True, "Keys work, but the response wasn't readable", "Got a 200 but couldn't parse the JSON."
        if isinstance(n, int) and n == 0:
            return True, "Keys work, but this index looks empty (0 cars)", \
                   "The App ID + Key are valid, but the Index Name may be wrong or the inventory feed is empty."
        return True, f"✅ Keys work — {n} vehicles in this index", \
               "This dealer is ready to resolve via the inventory API."
    if r.status_code in (401, 403):
        return False, "❌ Algolia rejected the key (403)", \
               "The saved key is rejected even with a matching Referer. First, re-grab a fresh key with the 🔑 bookmark in case the saved one is just stale. If a fresh key is ALSO rejected, the dealer's platform has locked it — it can't be reached from here, so mark it Ignore (not scannable)."
    if r.status_code == 404:
        return False, "❌ Index not found (404)", \
               f"The App ID + Key may be fine, but the Index Name '{index}' doesn't exist. Re-check the index."
    return False, f"❌ Algolia returned {r.status_code}", str(r.text)[:200]

# One-click bookmarklet: runs in the salesperson's OWN browser (no firewall can
# block it). It scrapes App ID + Index from the page AND hooks the network layer
# to capture the Search Key from a live Algolia request (many dealers only send
# the key as a request header when a search actually runs — not in the page HTML).
# Flow: click once to arm + scan. If the key isn't in the page, type a letter in
# the dealer's inventory search box, then click the bookmark again to grab it.
BOOKMARKLET_JS = (
    "javascript:(function(){try{"
    "var W=window;"
    "if(!W.__ffArmed){W.__ffArmed=1;"
    "W.__ffKeys={app:'',key:'',idx:'',dom:location.hostname.replace(/^www\\./,'')};"
    "var b=document.documentElement.innerHTML+'\\n'+Array.prototype.map.call(document.scripts,function(s){return s.textContent||''}).join('\\n');"
    "try{b+='\\n'+performance.getEntriesByType('resource').map(function(e){return e.name}).join('\\n')}catch(e){}"
    "function M(r){var x=b.match(r);return x?x[1]:''}"
    "W.__ffKeys.app=M(/([A-Z0-9]{10})-dsn\\.algolia\\.net/)||M(/(?:application-?id|appid|app_id)[\"'\\s:=]+([A-Z0-9]{10})/i);"
    "W.__ffKeys.idx=M(/[\"']?index(?:name)?[\"']?\\s*[:=]\\s*[\"']([a-z0-9_\\-]*inventory[a-z0-9_\\-]*)[\"']/i)||M(/([a-z0-9][a-z0-9_\\-]*_production_inventory[a-z0-9_]*)/i);"
    "var k0=M(/(?:x-algolia-api-key|api-?key|apikey|searchapikey)[\"'\\s:=]+[\"']?([A-Za-z0-9+\\/=_-]{20,})/i);if(k0)W.__ffKeys.key=k0;"
    "if(!W.__ffKeys.key){var ku=b.match(/x-algolia-api-key=([^&\\s\"']+)/i);if(ku)W.__ffKeys.key=decodeURIComponent(ku[1])}"
    "var oO=XMLHttpRequest.prototype.open,oS=XMLHttpRequest.prototype.setRequestHeader;"
    "XMLHttpRequest.prototype.open=function(m,u){try{if(/algolia/i.test(u)){var a=String(u).match(/([A-Z0-9]{10})-dsn/);if(a&&!W.__ffKeys.app)W.__ffKeys.app=a[1];var k=String(u).match(/x-algolia-api-key=([^&]+)/i);if(k)W.__ffKeys.key=decodeURIComponent(k[1])}}catch(e){}return oO.apply(this,arguments)};"
    "XMLHttpRequest.prototype.setRequestHeader=function(h,v){try{var H=String(h).toLowerCase();if(H=='x-algolia-api-key')W.__ffKeys.key=v;if(H=='x-algolia-application-id'&&v)W.__ffKeys.app=v}catch(e){}return oS.apply(this,arguments)};"
    "var oF=W.fetch;if(oF){W.fetch=function(res,opt){try{var u=(res&&res.url)||res;if(/algolia/i.test(String(u))){var a=String(u).match(/([A-Z0-9]{10})-dsn/);if(a&&!W.__ffKeys.app)W.__ffKeys.app=a[1];var k=String(u).match(/x-algolia-api-key=([^&]+)/i);if(k)W.__ffKeys.key=decodeURIComponent(k[1])}if(opt&&opt.headers){var H=opt.headers;var g=function(n){try{return H.get?H.get(n):(H[n]||H[n.toLowerCase()])}catch(e){return''}};var kk=g('x-algolia-api-key');if(kk)W.__ffKeys.key=kk;var ap=g('x-algolia-application-id');if(ap)W.__ffKeys.app=ap}}catch(e){}return oF.apply(this,arguments)}}"
    "}"
    "var F=W.__ffKeys;var o=[F.dom,F.app,F.key,F.idx].join('|');"
    "if(F.app&&F.key&&F.idx){"
    "if(navigator.clipboard&&navigator.clipboard.writeText){navigator.clipboard.writeText(o).then(function(){window.prompt('COPIED! Paste this into the Vault:',o)},function(){window.prompt('Copy this line into the Vault:',o)})}"
    "else{window.prompt('Copy this line into the Vault:',o)}"
    "}else{"
    "window.prompt('Almost! Got the App ID + Index, but not the Search Key yet.\\n\\nType any letter into THIS dealer\\'s inventory SEARCH box so a live search runs, then click this bookmark AGAIN.\\n\\n(Partial result - do not paste yet:)',o)"
    "}"
    "}catch(e){alert('Bookmarklet error: '+e.message)}})();"
)

# --- LOGIN ---
def check_password():
    if "password_correct" not in st.session_state:
        st.session_state.password_correct = False
    if st.session_state.password_correct:
        return True
    st.title("🔒 Auto-Analyst Login")
    password = st.text_input("Enter Company Password", type="password")
    if st.button("Log In"):
        if password == "tegna2026": 
            st.session_state.password_correct = True
            st.rerun()
        else:
            st.error("Incorrect password")
    return False

if not check_password():
    st.stop()

# --- SHARED EXPORT VISUALS (Premion palette) ---
# Charts are drawn with matplotlib so they render to PNG reliably in any
# environment (Plotly needs kaleido, which isn't guaranteed on the host). The
# data prep mirrors the on-screen dashboard exactly, so exports match the app.
PREMION_NAVY = "#0A1F5C"
PREMION_PERIWINKLE = "#8593E8"
PREMION_ORANGE = "#E8743B"
PREMION_BLUE = "#2E6FB7"
PREMION_GREEN = "#5FBFA0"
PREMION_RED = "#C0504D"

# Trim/drivetrain tokens used to roll vehicle names up to Make + Model for the
# aggregated "Top Sold/Missed Models" views. Detail tables keep the full name.
TRIM_TOKENS = {
    'lt', 'ltz', 'ls', 'rs', 'ss', 'rst', 'z71', 'zr2', 'premier', 'custom',
    'high', 'denali', 'elevation', 'slt', 'sle', 'at4', 'pro', 'limited',
    'laramie', 'big', 'horn', 'tradesman', 'rebel', 'trx', 'sport', 'touring',
    'ex', 'lx', 'se', 'sel', 'sv', 'sl', 's', 'xle', 'xse', 'sr', 'sr5', 'trd',
    'platinum', 'lariat', 'xlt', 'xl', 'king', 'ranch', 'raptor', 'tremor',
    'wildtrak', 'badlands', 'willys', 'rubicon', 'sahara', 'overland',
    'summit', 'trailhawk', 'latitude', 'altitude', 'fwd', 'awd', 'rwd', '4wd',
    '4x4', '4x2', '2wd', 'front', 'rear', 'four', 'all', 'wheel', 'drive',
    'quad', 'crew', 'double', 'extended', 'regular', 'cab', 'ultimate',
    'reserve', 'signature', 'premium', 'work', 'truck', 'trail', 'boss',
    'activ', 'redline', 'midnight', 'north', 'edition', 'hybrid', 'plug-in',
    'phev', 'ev', 'diesel', 'turbo', 'plus', 'select', 'preferred', 'essence',
    'avenir', 'st', 'st-line', 'titanium', 'calligraphy', 'ultimate,',
}

def normalize_model(name):
    """'2025 Chevrolet Silverado 1500 Rst Four Wheel...' -> 'Chevrolet Silverado 1500'."""
    s = re.sub(r'^\d{4}\s+', '', str(name)).strip()
    tokens = s.split()
    keep = []
    for i, t in enumerate(tokens):
        if i >= 1 and t.lower().strip('.,') in TRIM_TOKENS:
            break
        keep.append(t)
    if len(keep) < 2 and len(tokens) >= 2:
        keep = tokens[:2]
    return ' '.join(keep) if keep else s


def _fig_to_png_bytes(fig, w_in, h_in, tight=True):
    fig.set_size_inches(w_in, h_in)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150,
                bbox_inches=("tight" if tight else None),
                facecolor=fig.get_facecolor(), edgecolor="none")
    plt.close(fig)
    buf.seek(0)
    return buf

def _donut_chart(values, labels, colors, title, center_top, center_sub):
    """
    Fixed-geometry donut: the axes box is pinned so every donut renders the
    same circle size regardless of label/legend length (saved with tight=False).
    """
    fig, ax = plt.subplots(figsize=(4.0, 4.0))
    fig.patch.set_facecolor("white")
    ax.set_position([0.06, 0.20, 0.88, 0.66])
    total = sum(values)
    wedges, _texts, autotexts = ax.pie(
        values, colors=colors, startangle=90, counterclock=False,
        autopct=lambda p: f"{p:.0f}%" if p >= 5 else "",
        pctdistance=0.78,
        wedgeprops={'width': 0.42, 'edgecolor': 'white', 'linewidth': 2},
        textprops={'color': 'white', 'fontweight': 'bold', 'fontsize': 9})
    ax.text(0, 0.10, center_top, ha='center', va='center',
            fontsize=22, fontweight='bold', color=PREMION_NAVY)
    ax.text(0, -0.24, center_sub, ha='center', va='center',
            fontsize=9, color='#6B7280')
    fig.suptitle(title, y=0.965, fontweight='bold', color=PREMION_NAVY, fontsize=13)
    safe_labels = [str(l).replace('$', r'\$') for l in labels]  # avoid mathtext
    fig.legend(wedges, [f"{l}  ({v})" for l, v in zip(safe_labels, values)],
               loc='lower center', bbox_to_anchor=(0.5, 0.015),
               ncol=2, frameon=False, fontsize=8.5)
    return _fig_to_png_bytes(fig, 4.0, 4.0, tight=False)

def build_summary_chart_images(df, sold_df):
    """
    Render the three Executive Summary charts (Traffic Mix, Sales Mix, Sold Value
    Tiers) to PNG bytes for embedding in PDF/PPTX. Returns a dict of BytesIO
    (keys may be absent when there's no sold data). Never raises.
    """
    images = {}
    try:
        traffic = (df.groupby('Category')['Attributed Unique Visitors'].sum()
                     .reset_index().sort_values('Attributed Unique Visitors', ascending=True))
        fig, ax = plt.subplots()
        fig.patch.set_facecolor("white")
        colors = [PREMION_BLUE] * len(traffic)
        if colors:
            colors[-1] = PREMION_NAVY  # spotlight the top category
        bars = ax.barh(traffic['Category'], traffic['Attributed Unique Visitors'],
                       color=colors, height=0.62)
        ax.bar_label(bars, padding=4, fontsize=9, fontweight='bold', color=PREMION_NAVY)
        ax.set_title("Traffic Mix (Unique Visits by Page Type)", loc='left',
                     fontweight='bold', color=PREMION_NAVY, fontsize=13)
        ax.tick_params(axis='y', labelsize=9)
        ax.set_xticks([])
        for s in ('top', 'right', 'bottom'):
            ax.spines[s].set_visible(False)
        ax.spines['left'].set_color('#D1D5DB')
        max_v = traffic['Attributed Unique Visitors'].max()
        if pd.notna(max_v) and max_v > 0:
            ax.set_xlim(0, max_v * 1.15)
        ax.margins(y=0.02)
        images['traffic'] = _fig_to_png_bytes(fig, 6.6, 3.7)
    except Exception:
        pass
    if sold_df is not None and not sold_df.empty:
        try:
            tc = sold_df['Type'].value_counts()
            cmap = {'New': PREMION_BLUE, 'Used': PREMION_RED}
            images['salesmix'] = _donut_chart(
                tc.values.tolist(), tc.index.tolist(),
                [cmap.get(t, PREMION_PERIWINKLE) for t in tc.index],
                "Sales Mix (New vs Used)", str(int(tc.sum())), "Sold")
        except Exception:
            pass
        try:
            tier = sold_df['Price Tier'].value_counts()
            palette = [PREMION_GREEN, PREMION_ORANGE, PREMION_PERIWINKLE, PREMION_BLUE, PREMION_RED]
            images['tiers'] = _donut_chart(
                tier.values.tolist(), tier.index.tolist(),
                palette[:len(tier)],
                "Sold Value Tiers", str(int(tier.sum())), "Sold")
        except Exception:
            pass
    return images

def build_kpi_band_image(metrics):
    """
    Render the Executive Summary top-line as a navy dashboard-style KPI band
    (PNG bytes) for the PDF. Returns None on any failure so the PDF can fall
    back to plain text.
    """
    try:
        fig = plt.figure(figsize=(10.6, 1.6))
        fig.patch.set_facecolor(PREMION_NAVY)
        kpis = [
            ("UNITS SOLD (ATTRIBUTED)", str(metrics.get('units_sold', 0))),
            ("EST. REVENUE SOLD", f"${float(metrics.get('rev_sold', 0)):,.0f}"),
            ("TOTAL PIPELINE VALUE", f"${float(metrics.get('pipeline', 0)):,.0f}"),
            ("LOOK-TO-BOOK RATIO", f"{metrics.get('ltb', 0)}%"),
        ]
        for i, (lab, val) in enumerate(kpis):
            x = (i + 0.5) / 4.0
            fig.text(x, 0.72, lab, ha='center', va='center',
                     color='white', fontsize=8.5, fontweight='bold')
            fig.text(x, 0.36, val, ha='center', va='center',
                     color='white', fontsize=21, fontweight='bold')
        fig.text(3.5 / 4.0, 0.10,
                 f"New: {metrics.get('new_ltb', '-')}%  |  Used: {metrics.get('used_ltb', '-')}%",
                 ha='center', va='center', color='#DDE3F8', fontsize=8)
        import matplotlib.lines as mlines
        for i in (1, 2, 3):
            fig.add_artist(mlines.Line2D([i / 4.0, i / 4.0], [0.18, 0.84],
                                         transform=fig.transFigure,
                                         color=PREMION_PERIWINKLE, alpha=0.35, lw=1))
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150, facecolor=PREMION_NAVY)
        plt.close(fig)
        buf.seek(0)
        return buf
    except Exception:
        return None

def build_pptx_report(df, sold_df, metrics, chart_images, report_title,
                      missed_df=None, dealer_group=None,
                      include_dealer_details=False, include_missed=True):
    """
    Build the full Premion-styled PowerPoint deck: Executive Summary, charts,
    Top Sold, Missed Opportunities (Watch List), Auto Group scoreboard +
    per-dealer deep dives for multi-dealer reports, and Glossary/Methodology.
    Returns PPTX file bytes, or None if python-pptx isn't available.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
    except Exception:
        return None

    NAVY = RGBColor(0x0A, 0x1F, 0x5C)
    PERI = RGBColor(0x85, 0x93, 0xE8)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x33, 0x33, 0x33)
    GREY = RGBColor(0x6B, 0x72, 0x80)
    LIGHT = RGBColor(0xEF, 0xF1, 0xFA)

    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_text(slide, left, top, width, height, text, size, color,
                 bold=False, align=PP_ALIGN.LEFT, italic=False):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
        r.font.italic = italic
        r.font.name = "Arial"
        return tb

    def fill_rect(slide, left, top, width, height, color, rounded=True):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        shp = slide.shapes.add_shape(shape_type, left, top, width, height)
        shp.fill.solid(); shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        return shp

    def add_footer(slide):
        add_text(slide, Inches(0.5), Inches(7.02), Inches(6), Inches(0.35),
                 "TEGNA  |  PREMION", 12, NAVY, bold=True)

    def new_slide(title, band_h=1.1):
        s = prs.slides.add_slide(blank)
        fill_rect(s, Inches(0), Inches(0), prs.slide_width, Inches(band_h), NAVY,
                  rounded=False)
        add_text(s, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7),
                 title, 26, PERI, bold=True)
        add_footer(s)
        return s

    def add_table(slide, left, top, width, rows_data, col_fracs, font_size=10,
                  row_h_in=0.32):
        if not rows_data:
            return None
        n_rows = len(rows_data); n_cols = len(rows_data[0])
        gfx = slide.shapes.add_table(n_rows, n_cols, left, top, width,
                                     Inches(row_h_in * n_rows))
        table = gfx.table
        total = float(sum(col_fracs))
        for j, frac in enumerate(col_fracs):
            table.columns[j].width = int(width * (frac / total))
        for i, row in enumerate(rows_data):
            for j, val in enumerate(row):
                cell = table.cell(i, j)
                cell.margin_left = Pt(4); cell.margin_right = Pt(4)
                cell.margin_top = Pt(1); cell.margin_bottom = Pt(1)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf = cell.text_frame; tf.word_wrap = False
                p = tf.paragraphs[0]
                r = p.add_run(); r.text = str(val)
                r.font.name = "Arial"; r.font.size = Pt(font_size)
                if i == 0:
                    r.font.bold = True; r.font.color.rgb = WHITE
                    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
                else:
                    r.font.color.rgb = DARK
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT if i % 2 == 0 else WHITE
        return gfx

    def add_picture_fit(slide, buf, left, top, max_w_in, max_h_in):
        """Place an image inside a bounding box, preserving aspect ratio."""
        try:
            from PIL import Image as PILImage
            buf.seek(0)
            with PILImage.open(buf) as im:
                aspect = im.width / im.height
        except Exception:
            aspect = 1.0
        buf.seek(0)
        w = max_w_in; h = w / aspect
        if h > max_h_in:
            h = max_h_in; w = h * aspect
        pic = slide.shapes.add_picture(buf, Inches(left), Inches(top),
                                       width=Inches(w))
        buf.seek(0)
        return pic

    def _model_counts(frame, count_label):
        f = frame.copy()
        f['Model_Only'] = f['Vehicle Name'].apply(normalize_model)
        counts = f['Model_Only'].value_counts().reset_index()
        counts.columns = ['Make/Model', count_label]
        return counts[counts[count_label] > 1].head(10)

    def _fmt_money(v):
        try:
            return f"${float(v):,.0f}"
        except Exception:
            return str(v)

    missed_df = missed_df if missed_df is not None else pd.DataFrame()
    has_missed = include_missed and not missed_df.empty
    has_group = dealer_group is not None and not dealer_group.empty

    # --- Slide 1: title band + KPI cards ---
    s1 = prs.slides.add_slide(blank)
    fill_rect(s1, Inches(0), Inches(0), prs.slide_width, Inches(1.6), NAVY,
              rounded=False)
    add_text(s1, Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.8),
             "Premion Website Attribution", 34, PERI, bold=True)
    # Client-facing: show the report DATE, not the pull time.
    display_title = re.sub(r'\s*\(\d{1,2}:\d{2}\s*[AP]M[^)]*\)', '', str(report_title))
    display_title = re.sub(r'\s*-?\s*\d{1,2}:\d{2}\s*[AP]M(\s*ET)?\s*$', '', display_title).strip()
    try:
        report_date = datetime.datetime.now(pytz.timezone("US/Eastern")).strftime("%B %d, %Y").replace(" 0", " ")
    except Exception:
        report_date = datetime.datetime.now().strftime("%B %d, %Y")
    add_text(s1, Inches(0.5), Inches(0.9), Inches(12.3), Inches(0.5),
             f"Inventory Impact Report  |  {display_title}  |  {report_date}",
             18, WHITE, bold=True)
    add_text(s1, Inches(0.5), Inches(1.75), Inches(12.3), Inches(0.5),
             "Complementary to your website attribution: the vehicles our audience "
             "shopped that have since moved off the lot.", 13, GREY, italic=True)

    units_val = metrics.get('units_sold', metrics.get('units', 0))
    rev_val = metrics.get('rev_sold', metrics.get('revenue', 0))
    ltb_val = metrics.get('ltb', 0)
    kpis = [("Units Sold (Attributed)", str(units_val)),
            ("Est. Revenue Sold", _fmt_money(rev_val)),
            ("Total Pipeline Value", _fmt_money(metrics.get('pipeline', 0))),
            ("Look-to-Book Ratio", f"{ltb_val}%")]
    card_w = Inches(2.9); card_h = Inches(1.9); gap = Inches(0.28)
    total_w = card_w * 4 + gap * 3
    x0 = (prs.slide_width - total_w) / 2
    y0 = Inches(2.7)
    for i, (label, val) in enumerate(kpis):
        cx = x0 + i * (card_w + gap)
        fill_rect(s1, cx, y0, card_w, card_h, NAVY)
        add_text(s1, cx, y0 + Inches(0.25), card_w, Inches(0.5), label, 12, WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(s1, cx, y0 + Inches(0.8), card_w, Inches(0.9), val, 30, WHITE,
                 bold=True, align=PP_ALIGN.CENTER)
    if metrics.get('new_ltb') is not None or metrics.get('used_ltb') is not None:
        add_text(s1, x0 + 3 * (card_w + gap), y0 + card_h + Inches(0.08),
                 card_w, Inches(0.4),
                 f"New: {metrics.get('new_ltb', '-')}%  |  Used: {metrics.get('used_ltb', '-')}%",
                 11, GREY, align=PP_ALIGN.CENTER)
    add_text(s1, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.5),
             "Look-to-Book = Sold VDPs / Total Active VDPs - the conversion "
             "velocity of the inventory we advertised.", 12, GREY, italic=True,
             align=PP_ALIGN.CENTER)
    add_footer(s1)

    # --- Slide 2: Traffic Mix (full-size) ---
    if chart_images.get('traffic'):
        s2 = new_slide("Traffic Mix: Where the Audience Shopped")
        add_picture_fit(s2, chart_images['traffic'], 1.5, 1.5, 10.3, 4.9)
        add_text(s2, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.5),
                 "VDP visits are high-intent inventory shoppers - the audience "
                 "most likely to buy.", 12, GREY, italic=True,
                 align=PP_ALIGN.CENTER)

    # --- Slide 3: Sales Mix & Sold Value Tiers ---
    if chart_images.get('salesmix') or chart_images.get('tiers'):
        s2b = new_slide("What Sold: Sales Mix & Value Tiers")
        if chart_images.get('salesmix') and chart_images.get('tiers'):
            add_picture_fit(s2b, chart_images['salesmix'], 1.55, 1.55, 4.6, 4.6)
            add_picture_fit(s2b, chart_images['tiers'], 7.2, 1.55, 4.6, 4.6)
        else:
            only = chart_images.get('salesmix') or chart_images.get('tiers')
            add_picture_fit(s2b, only, 4.35, 1.55, 4.6, 4.6)
        add_text(s2b, Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.5),
                 "The mix of inventory our traffic helped move - what sold, and "
                 "at what price points.", 12, GREY, italic=True,
                 align=PP_ALIGN.CENTER)

    # --- Slide 3: Top Sold ---
    if sold_df is not None and not sold_df.empty:
        s3 = new_slide("Top Sold: Proof of Inventory Movement")
        top_models = _model_counts(sold_df, 'Units Sold')
        if not top_models.empty:
            add_text(s3, Inches(0.5), Inches(1.35), Inches(4.6), Inches(0.4),
                     "Top Sold Models (>1 Unit)", 14, NAVY, bold=True)
            rows = [["Make/Model", "Units"]] + [
                [str(r['Make/Model'])[:30], str(r['Units Sold'])]
                for _, r in top_models.iterrows()]
            add_table(s3, Inches(0.5), Inches(1.8), Inches(4.4), rows,
                      [4, 1], font_size=10)
        detail_x = 5.3 if not top_models.empty else 0.5
        detail_w = 7.5 if not top_models.empty else 12.3
        add_text(s3, Inches(detail_x), Inches(1.35), Inches(detail_w), Inches(0.4),
                 "Top Sold Units (Detail)", 14, NAVY, bold=True)
        top_sold = sold_df.sort_values('Attributed Unique Visitors', ascending=False).head(10)
        rows = [["Vehicle", "Type", "Visitors", "VIN"]] + [
            [str(r['Vehicle Name'])[:34], str(r['Type']),
             str(r['Attributed Unique Visitors']), str(r['VIN'])]
            for _, r in top_sold.iterrows()]
        add_table(s3, Inches(detail_x), Inches(1.8), Inches(detail_w), rows,
                  [3.4, 0.9, 1.0, 2.4], font_size=9)

    # --- Slide 4: Missed Opportunities ---
    if has_missed:
        s4 = new_slide("Missed Opportunities: The Watch List")
        add_text(s4, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.6),
                 "Active vehicles receiving above-average campaign traffic that "
                 "haven't sold yet - review these VDPs for missing photos, "
                 "\"Call for Price\" buttons, or pricing outliers.", 12, GREY,
                 italic=True)
        top_missed = missed_df.sort_values('Attributed Unique Visitors', ascending=False).head(10)
        rows = [["Vehicle", "Type", "Visitors", "VIN"]] + [
            [str(r['Vehicle Name'])[:40], str(r['Type']),
             str(r['Attributed Unique Visitors']), str(r['VIN'])]
            for _, r in top_missed.iterrows()]
        add_table(s4, Inches(0.5), Inches(2.0), Inches(12.3), rows,
                  [4.2, 0.9, 1.0, 2.4], font_size=10)

    # --- Auto Group scoreboard (multi-dealer) ---
    if has_group:
        group_rows = dealer_group.reset_index(drop=True)
        per_slide = 12
        n_pages = (len(group_rows) + per_slide - 1) // per_slide
        for page in range(n_pages):
            title = "Auto Group Scoreboard"
            if n_pages > 1:
                title += f" ({page + 1} of {n_pages})"
            sg = new_slide(title)
            chunk = group_rows.iloc[page * per_slide:(page + 1) * per_slide]
            rows = [["Dealer", "Traffic", "VDPs", "Sold", "LTB", "Est. Rev Sold", "Pipeline"]]
            for _, r in chunk.iterrows():
                rows.append([str(r['Dealer'])[:32], str(r['Total Visitors']),
                             str(r['VDPs Shopped']), str(r['Units Sold']),
                             f"{r['Look-to-Book (%)']}%",
                             _fmt_money(r['Est. Rev Sold']),
                             _fmt_money(r['Pipeline Value'])])
            add_table(sg, Inches(0.5), Inches(1.5), Inches(12.3), rows,
                      [3.4, 1.0, 0.8, 0.8, 0.9, 1.5, 1.5], font_size=10)

    # --- Dealer Deep Dives (one slide per dealer) ---
    if include_dealer_details and has_group:
        for _, grow in dealer_group.iterrows():
            dealer = grow['Dealer']
            d_sold = sold_df[sold_df['Dealer'] == dealer] if sold_df is not None and not sold_df.empty else pd.DataFrame()
            d_missed = missed_df[missed_df['Dealer'] == dealer] if not missed_df.empty else pd.DataFrame()
            if d_sold.empty and d_missed.empty:
                continue
            sd = new_slide(f"Dealer Profile: {str(dealer)[:48]}")
            add_text(sd, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.4),
                     f"Traffic: {grow['Total Visitors']}   |   VDPs Shopped: {grow['VDPs Shopped']}   |   "
                     f"Units Sold: {grow['Units Sold']}   |   LTB: {grow['Look-to-Book (%)']}%   |   "
                     f"Est. Rev Sold: {_fmt_money(grow['Est. Rev Sold'])}   |   "
                     f"Pipeline: {_fmt_money(grow['Pipeline Value'])}",
                     12, NAVY, bold=True)
            if not d_sold.empty:
                add_text(sd, Inches(0.5), Inches(1.85), Inches(5.9), Inches(0.4),
                         "Top Sold Units", 13, NAVY, bold=True)
                d_top = d_sold.sort_values('Attributed Unique Visitors', ascending=False).head(5)
                rows = [["Vehicle", "Type", "Visitors"]] + [
                    [str(r['Vehicle Name'])[:34], str(r['Type']),
                     str(r['Attributed Unique Visitors'])]
                    for _, r in d_top.iterrows()]
                add_table(sd, Inches(0.5), Inches(2.25), Inches(5.9), rows,
                          [3.6, 0.9, 1.0], font_size=9)
            if not d_missed.empty:
                add_text(sd, Inches(6.9), Inches(1.85), Inches(5.9), Inches(0.4),
                         "Missed Opportunities (Watch List)", 13, NAVY, bold=True)
                d_topm = d_missed.sort_values('Attributed Unique Visitors', ascending=False).head(5)
                rows = [["Vehicle", "Type", "Visitors"]] + [
                    [str(r['Vehicle Name'])[:34], str(r['Type']),
                     str(r['Attributed Unique Visitors'])]
                    for _, r in d_topm.iterrows()]
                add_table(sd, Inches(6.9), Inches(2.25), Inches(5.9), rows,
                          [3.6, 0.9, 1.0], font_size=9)

    # --- Glossary & Methodology ---
    sg2 = new_slide("Glossary & Methodology")
    glossary = [
        ("Units Sold (Attributed)",
         "Vehicles removed from the dealer's live inventory after receiving attributed campaign traffic."),
        ("Est. Revenue Sold / Pipeline Value",
         "Directional estimates. New: base MSRP for the model. Used: base MSRP depreciated by age "
         "(15% year 1, 10% each following year). Not exact transaction prices."),
        ("Look-to-Book Ratio",
         "Sold VDPs / Total Active VDPs - the conversion velocity of the inventory, split New vs Used."),
        ("Traffic Mix",
         "Where audiences navigated on the site: VDPs, Service, Search, Incentives/Offers, Homepage."),
        ("Missed Opportunities (Watch List)",
         "Active vehicles with above-average traffic that haven't sold - flags potential pricing or "
         "merchandising issues."),
        ("Methodology Note",
         "Inventory status reflects the dealer's website at the moment this report was run, not a "
         "historical snapshot."),
    ]
    tb = sg2.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.2))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    from pptx.util import Pt as _Pt
    for head, body in glossary:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = head
        r.font.name = "Arial"; r.font.size = _Pt(14); r.font.bold = True
        r.font.color.rgb = NAVY
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = body
        r2.font.name = "Arial"; r2.font.size = _Pt(11); r2.font.color.rgb = DARK
        p2.space_after = _Pt(8)

    out = io.BytesIO()
    prs.save(out); out.seek(0)
    return out.read()

# --- PDF GENERATOR FUNCTION ---
def create_pdf_report(df, sold_df, metrics, missed_df, include_missed, dealer_group=None, include_dealer_details=False, chart_images=None):
    # HELPER: Sanitizes scraped text so it doesn't crash the PDF encoder
    def safe_str(text):
        if pd.isna(text): return ""
        text = str(text)
        reps = {'“':'"', '”':'"', '‘':"'", '’':"'", '—':'-', '–':'-', '®':'', '™':'', '©':''}
        for k, v in reps.items(): text = text.replace(k, v)
        return text.encode('latin-1', 'ignore').decode('latin-1')

    # HELPER: navy header row for tables (sleeker look, matches the KPI band)
    def thead(cols, size=9):
        pdf.set_font("Arial", "B", size)
        pdf.set_fill_color(10, 31, 92)
        pdf.set_text_color(255, 255, 255)
        for w, t in cols:
            pdf.cell(w, 8, t, border=1, fill=True)
        pdf.ln()
        pdf.set_text_color(0, 0, 0)
        pdf.set_fill_color(240, 240, 240)

    pdf = FPDF()
    pdf.add_page()
    
    eastern = pytz.timezone('US/Eastern')
    current_time = datetime.datetime.now(eastern)
    
    domain_count = df['Dealer'].nunique()
    report_title = "Auto-Sales Intelligence Report"
    if domain_count == 1:
        report_title = f"{df['Dealer'].iloc[0]} Intelligence Report"

    pdf.set_font("Arial", "B", 18)
    pdf.cell(0, 15, safe_str(report_title), ln=True, align="C")
    
    pdf.set_font("Arial", "I", 10)
    pdf.cell(0, 8, f"Generated on: {current_time.strftime('%Y-%m-%d %I:%M %p ET')} | Attribution Threshold: {metrics.get('min_visitors', 1)}+ VDP Visitors", ln=True, align="C")
    pdf.ln(5)

    pdf.set_font("Arial", "B", 14)
    pdf.set_fill_color(240, 240, 240)
    pdf.cell(0, 10, " 1. Executive Summary", ln=True, fill=True)
    pdf.ln(4)

    # Navy dashboard-style KPI band; falls back to plain text if imaging fails.
    band_done = False
    kpi_band = build_kpi_band_image(metrics)
    if kpi_band is not None:
        try:
            import tempfile
            tfb = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
            tfb.write(kpi_band.read()); tfb.close()
            pdf.image(tfb.name, x=15, w=180)
            try: os.remove(tfb.name)
            except Exception: pass
            band_done = True
            pdf.ln(2)
            pdf.set_font("Arial", "I", 8)
            pdf.set_text_color(90, 90, 90)
            pdf.cell(0, 5, "Look-to-Book = Sold VDPs / Total Active VDPs - the conversion velocity of the inventory we advertised.", ln=True, align="C")
            pdf.set_text_color(0, 0, 0)
        except Exception:
            band_done = False
    if not band_done:
        pdf.set_font("Arial", "", 12)
        col_width = pdf.w / 2.2
        pdf.cell(col_width, 8, f"Total Units Sold: {metrics['units_sold']}", border=0)
        pdf.cell(col_width, 8, f"Est. Revenue Sold: ${metrics['rev_sold']:,.0f}", border=0, ln=True)
        pdf.cell(col_width, 8, f"Pipeline Value: ${metrics['pipeline']:,.0f}", border=0)
        pdf.cell(col_width, 8, f"Look-to-Book Ratio: {metrics['ltb']}%", border=0, ln=True)
        pdf.set_font("Arial", "I", 10)
        pdf.cell(col_width, 6, "", border=0)
        pdf.cell(col_width, 6, f"(New: {metrics['new_ltb']}% | Used: {metrics['used_ltb']}%)", border=0, ln=True)
    pdf.ln(8)

    section_offset = 0
    if dealer_group is not None and not dealer_group.empty:
        pdf.set_font("Arial", "B", 14)
        pdf.cell(0, 10, " 2. Auto Group / Tier 2 Breakdown", ln=True, fill=True)
        pdf.ln(5)
        
        thead([(50, "Dealer Name"), (18, "Traffic"), (15, "VDPs"), (15, "Sold"),
               (15, "LTB"), (35, "Est. Rev Sold"), (35, "Pipeline Value")])
        
        pdf.set_font("Arial", "", 8)
        for _, row in dealer_group.head(15).iterrows():
            pdf.cell(50, 8, safe_str(row['Dealer'])[:28], border=1)
            pdf.cell(18, 8, str(row['Total Visitors']), border=1)
            pdf.cell(15, 8, str(row['VDPs Shopped']), border=1)
            pdf.cell(15, 8, str(row['Units Sold']), border=1)
            pdf.cell(15, 8, f"{row['Look-to-Book (%)']}%", border=1)
            pdf.cell(35, 8, f"${row['Est. Rev Sold']:,.0f}", border=1)
            pdf.cell(35, 8, f"${row['Pipeline Value']:,.0f}", border=1)
            pdf.ln()
        pdf.ln(8)
        section_offset = 1

    sec_market = 2 + section_offset
    pdf.set_font("Arial", "B", 14)
    pdf.cell(0, 10, f" {sec_market}. Market Insights", ln=True, fill=True)
    pdf.ln(5)

    # Visual charts (matching the on-screen dashboard) embedded as images.
    # Heights are computed from the real PNG dimensions and page breaks are
    # inserted before placing, so charts can never run off the page edge.
    if chart_images:
        try:
            import tempfile
            try:
                from PIL import Image as PILImage
            except Exception:
                PILImage = None

            def _img_h_mm(path, w_mm):
                if PILImage is not None:
                    try:
                        with PILImage.open(path) as im:
                            return w_mm * im.height / im.width
                    except Exception:
                        pass
                return w_mm * 0.62  # safe fallback aspect

            def _ensure_space(h_mm):
                if pdf.get_y() + h_mm > pdf.h - 15:
                    pdf.add_page()

            def _caption(text):
                pdf.set_font("Arial", "I", 8)
                pdf.set_text_color(90, 90, 90)
                pdf.cell(0, 5, safe_str(text), ln=True, align="C")
                pdf.set_text_color(0, 0, 0)

            saved = {}
            for key in ('traffic', 'salesmix', 'tiers'):
                if chart_images.get(key):
                    chart_images[key].seek(0)
                    tf = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
                    tf.write(chart_images[key].read()); tf.close()
                    saved[key] = tf.name
                    chart_images[key].seek(0)
            if saved.get('traffic'):
                h = _img_h_mm(saved['traffic'], 180)
                _ensure_space(h + 8)
                pdf.image(saved['traffic'], x=15, w=180)
                _caption("Where our audience navigated on the dealer's site. VDP visits = high-intent inventory shoppers.")
                pdf.ln(2)
            if saved.get('salesmix') or saved.get('tiers'):
                h1 = _img_h_mm(saved['salesmix'], 88) if saved.get('salesmix') else 0
                h2 = _img_h_mm(saved['tiers'], 88) if saved.get('tiers') else 0
                row_h = max(h1, h2)
                _ensure_space(row_h + 8)
                y = pdf.get_y()
                if saved.get('salesmix'):
                    pdf.image(saved['salesmix'], x=15, y=y, w=88)
                if saved.get('tiers'):
                    pdf.image(saved['tiers'], x=107, y=y, w=88)
                pdf.set_y(y + row_h + 1)
                _caption("What sold, and at what price points - the mix of inventory our traffic helped move.")
            for path in saved.values():
                try: os.remove(path)
                except Exception: pass
            pdf.ln(4)
        except Exception:
            pass  # charts are a bonus; never let them break the PDF

    pdf.set_font("Arial", "B", 11)
    pdf.cell(0, 8, "Traffic Mix (Unique Visits by Page Type)", ln=True)
    thead([(100, "Page Category"), (40, "Unique Visits")])
    pdf.set_font("Arial", "", 9)
    traffic_mix = df.groupby('Category')['Attributed Unique Visitors'].sum().reset_index().sort_values('Attributed Unique Visitors', ascending=False)
    for _, row in traffic_mix.iterrows():
        pdf.cell(100, 8, safe_str(row['Category']), border=1)
        pdf.cell(40, 8, str(row['Attributed Unique Visitors']), border=1)
        pdf.ln()
    pdf.ln(5)

    if not sold_df.empty:
        pdf.set_font("Arial", "B", 11)
        pdf.cell(0, 8, "Sales Mix (New vs Used)", ln=True)
        thead([(70, "Type"), (35, "Sold Count"), (35, "Share %")])
        pdf.set_font("Arial", "", 9)
        sales_mix = sold_df['Type'].value_counts(normalize=True).mul(100).round(1).reset_index()
        sales_mix.columns = ['Type', 'Share']
        sales_counts = sold_df['Type'].value_counts().reset_index()
        sales_counts.columns = ['Type', 'Count']
        merged_sales = pd.merge(sales_mix, sales_counts, on='Type')
        for _, row in merged_sales.iterrows():
            pdf.cell(70, 8, safe_str(row['Type']), border=1)
            pdf.cell(35, 8, str(row['Count']), border=1)
            pdf.cell(35, 8, f"{row['Share']}%", border=1)
            pdf.ln()
        pdf.ln(5)

    if not sold_df.empty:
        pdf.set_font("Arial", "B", 11)
        pdf.cell(0, 8, "Price Tiers (Sold Units)", ln=True)
        thead([(70, "Price Tier"), (35, "Sold Count"), (35, "Share %")])
        pdf.set_font("Arial", "", 9)
        tier_mix = sold_df['Price Tier'].value_counts(normalize=True).mul(100).round(1).reset_index()
        tier_mix.columns = ['Price Tier', 'Share']
        tier_counts = sold_df['Price Tier'].value_counts().reset_index()
        tier_counts.columns = ['Price Tier', 'Count']
        merged_tiers = pd.merge(tier_mix, tier_counts, on='Price Tier')
        for _, row in merged_tiers.iterrows():
            pdf.cell(70, 8, safe_str(row['Price Tier']), border=1)
            pdf.cell(35, 8, str(row['Count']), border=1)
            pdf.cell(35, 8, f"{row['Share']}%", border=1)
            pdf.ln()
    pdf.ln(10)

    if not sold_df.empty:
        sold_models = sold_df.copy()
        sold_models['Model_Only'] = sold_models['Vehicle Name'].apply(normalize_model)
        model_counts = sold_models['Model_Only'].value_counts().reset_index()
        model_counts.columns = ['Make/Model', 'Units Sold']
        top_models = model_counts[model_counts['Units Sold'] > 1].head(10)
        
        if not top_models.empty:
            sec_models = 3 + section_offset
            pdf.set_font("Arial", "B", 14)
            pdf.cell(0, 10, f" {sec_models}. Top Sold Models (Aggregated > 1 Unit)", ln=True, fill=True)
            pdf.ln(5)
            thead([(100, "Make/Model"), (40, "Units Sold")], size=10)
            pdf.set_font("Arial", "", 9)
            for _, row in top_models.iterrows():
                pdf.cell(100, 8, safe_str(row['Make/Model']), border=1)
                pdf.cell(40, 8, str(row['Units Sold']), border=1)
                pdf.ln()
            pdf.ln(5)

    pdf.set_font("Arial", "B", 11)
    pdf.cell(0, 10, "Top Sold Units (Detail)", ln=True)
    thead([(80, "Vehicle Name"), (25, "Type"), (20, "Visitors"), (65, "VIN")], size=10)

    if not sold_df.empty:
        top_sold = sold_df.sort_values('Attributed Unique Visitors', ascending=False).head(10)
        for _, row in top_sold.iterrows():
            name = safe_str(row['Vehicle Name'])[:35]
            pdf.set_font("Arial", "", 9) 
            pdf.cell(80, 8, name, border=1)
            pdf.cell(25, 8, safe_str(row['Type']), border=1)
            pdf.cell(20, 8, str(row['Attributed Unique Visitors']), border=1)
            pdf.set_font("Arial", "", 8)
            pdf.cell(65, 8, safe_str(row['VIN']), border=1)
            pdf.ln()
    else:
        pdf.set_font("Arial", "", 9)
        pdf.cell(0, 8, "No sales identified.", border=1, ln=True)

    if include_missed:
        pdf.add_page()
        sec_missed = 4 + section_offset
        pdf.set_font("Arial", "B", 14)
        pdf.cell(0, 10, f" {sec_missed}. Missed Opportunities (The Watch List)", ln=True, fill=True)
        pdf.ln(5)
        
        if not missed_df.empty:
            missed_models = missed_df.copy()
            missed_models['Model_Only'] = missed_models['Vehicle Name'].apply(normalize_model)
            missed_counts = missed_models['Model_Only'].value_counts().reset_index()
            missed_counts.columns = ['Make/Model', 'Count']
            top_missed = missed_counts[missed_counts['Count'] > 1].head(10)
            
            if not top_missed.empty:
                pdf.set_font("Arial", "B", 11)
                pdf.cell(0, 10, "Top Missed Models (Aggregated > 1 Unit)", ln=True)
                thead([(100, "Make/Model"), (40, "Missed Count")], size=10)
                pdf.set_font("Arial", "", 9)
                for _, row in top_missed.iterrows():
                    pdf.cell(100, 8, safe_str(row['Make/Model']), border=1)
                    pdf.cell(40, 8, str(row['Count']), border=1)
                    pdf.ln()
                pdf.ln(5)
        
        pdf.set_font("Arial", "B", 11)
        pdf.cell(0, 10, "Missed Opportunities (Detail)", ln=True)
        thead([(85, "Vehicle Name (Clickable)"), (65, "VIN"), (20, "Type"), (20, "Visitors")], size=10)
        pdf.set_font("Arial", "", 9)
        if not missed_df.empty:
             top_missed_detail = missed_df.sort_values('Attributed Unique Visitors', ascending=False).head(10)
             for _, row in top_missed_detail.iterrows():
                 name = safe_str(row['Vehicle Name'])[:35]
                 url = str(row['Page Url'])
                 pdf.set_text_color(0, 0, 255) 
                 pdf.cell(85, 8, name, border=1, link=url)
                 pdf.set_text_color(0, 0, 0)
                 pdf.set_font("Arial", "", 8) 
                 pdf.cell(65, 8, safe_str(row['VIN']), border=1)
                 pdf.set_font("Arial", "", 9)
                 pdf.cell(20, 8, safe_str(row['Type']), border=1)
                 pdf.cell(20, 8, str(row['Attributed Unique Visitors']), border=1)
                 pdf.ln()
                 
    # --- DEALER DEEP DIVES EXTENSION IN PDF ---
    if include_dealer_details and dealer_group is not None and not dealer_group.empty:
        for dealer in dealer_group['Dealer'].tolist():
            dealer_sold = sold_df[sold_df['Dealer'] == dealer]
            dealer_missed = missed_df[missed_df['Dealer'] == dealer]
            
            if dealer_sold.empty and dealer_missed.empty:
                continue 
                
            pdf.add_page()
            pdf.set_font("Arial", "B", 14)
            pdf.set_fill_color(220, 220, 220)
            pdf.cell(0, 10, f" Dealership Profile: {safe_str(dealer)}", ln=True, fill=True)
            pdf.ln(5)
            
            if not dealer_sold.empty:
                d_sold_models = dealer_sold.copy()
                d_sold_models['Model_Only'] = d_sold_models['Vehicle Name'].apply(normalize_model)
                d_model_counts = d_sold_models['Model_Only'].value_counts().reset_index()
                d_model_counts.columns = ['Make/Model', 'Units Sold']
                d_top_models = d_model_counts[d_model_counts['Units Sold'] > 1].head(5)
                
                if not d_top_models.empty:
                    pdf.set_font("Arial", "B", 11)
                    pdf.cell(0, 10, "Top Sold Models", ln=True)
                    thead([(100, "Make/Model"), (40, "Units Sold")], size=10)
                    pdf.set_font("Arial", "", 9)
                    for _, row in d_top_models.iterrows():
                        pdf.cell(100, 8, safe_str(row['Make/Model']), border=1)
                        pdf.cell(40, 8, str(row['Units Sold']), border=1)
                        pdf.ln()
                    pdf.ln(5)
                    
            if not dealer_missed.empty:
                d_missed_models = dealer_missed.copy()
                d_missed_models['Model_Only'] = d_missed_models['Vehicle Name'].apply(normalize_model)
                d_missed_counts = d_missed_models['Model_Only'].value_counts().reset_index()
                d_missed_counts.columns = ['Make/Model', 'Count']
                d_top_missed = d_missed_counts[d_missed_counts['Count'] > 1].head(5)
                
                if not d_top_missed.empty:
                    pdf.set_font("Arial", "B", 11)
                    pdf.cell(0, 10, "Top Missed Models", ln=True)
                    thead([(100, "Make/Model"), (40, "Missed Count")], size=10)
                    pdf.set_font("Arial", "", 9)
                    for _, row in d_top_missed.iterrows():
                        pdf.cell(100, 8, safe_str(row['Make/Model']), border=1)
                        pdf.cell(40, 8, str(row['Count']), border=1)
                        pdf.ln()
                    pdf.ln(5)

            if not dealer_sold.empty:
                pdf.set_font("Arial", "B", 11)
                pdf.cell(0, 10, "Top Sold Units (Detail)", ln=True)
                thead([(80, "Vehicle Name"), (25, "Type"), (20, "Visitors"), (65, "VIN")], size=10)
                d_top_sold = dealer_sold.sort_values('Attributed Unique Visitors', ascending=False).head(10)
                for _, row in d_top_sold.iterrows():
                    name = safe_str(row['Vehicle Name'])[:35]
                    pdf.set_font("Arial", "", 9) 
                    pdf.cell(80, 8, name, border=1)
                    pdf.cell(25, 8, safe_str(row['Type']), border=1)
                    pdf.cell(20, 8, str(row['Attributed Unique Visitors']), border=1)
                    pdf.set_font("Arial", "", 8)
                    pdf.cell(65, 8, safe_str(row['VIN']), border=1)
                    pdf.ln()
                pdf.ln(5)

            if not dealer_missed.empty:
                pdf.set_font("Arial", "B", 11)
                pdf.cell(0, 10, "Missed Opportunities (Detail)", ln=True)
                thead([(85, "Vehicle Name (Clickable)"), (65, "VIN"), (20, "Type"), (20, "Visitors")], size=10)
                d_top_missed_detail = dealer_missed.sort_values('Attributed Unique Visitors', ascending=False).head(10)
                pdf.set_font("Arial", "", 9)
                for _, row in d_top_missed_detail.iterrows():
                     name = safe_str(row['Vehicle Name'])[:35]
                     url = str(row['Page Url'])
                     pdf.set_text_color(0, 0, 255) 
                     pdf.cell(85, 8, name, border=1, link=url)
                     pdf.set_text_color(0, 0, 0)
                     pdf.set_font("Arial", "", 8) 
                     pdf.cell(65, 8, safe_str(row['VIN']), border=1)
                     pdf.set_font("Arial", "", 9)
                     pdf.cell(20, 8, safe_str(row['Type']), border=1)
                     pdf.cell(20, 8, str(row['Attributed Unique Visitors']), border=1)
                     pdf.ln()
                 
    # --- GLOSSARY & METHODOLOGY ---
    pdf.add_page()
    pdf.set_font("Arial", "B", 14)
    pdf.set_fill_color(240, 240, 240)
    pdf.cell(0, 10, " Glossary & Methodology", ln=True, fill=True)
    pdf.ln(4)
    glossary_items = [
        ("Units Sold (Attributed)",
         "Vehicles removed from the dealer's live inventory after receiving attributed campaign traffic."),
        ("Est. Revenue Sold / Pipeline Value",
         "Directional value estimates. New: base MSRP for the model. Used: base MSRP depreciated by age (15% year 1, 10% each following year). Not exact transaction prices; excludes trims, options, and dealer markups."),
        ("Look-to-Book Ratio",
         "Sold VDPs / Total Active VDPs - the conversion velocity of the inventory, split by New vs Used."),
        ("Traffic Mix",
         "Where audiences navigated on the site: VDPs, Service, Search, Incentives/Offers, Homepage."),
        ("Missed Opportunities (The Watch List)",
         "Active vehicles receiving above-average traffic that haven't sold yet. Review these VDPs for missing photos, 'Call for Price' buttons, or pricing outliers."),
        ("Methodology Note",
         "Inventory status reflects the dealer's website at the moment this report was run, not a historical snapshot."),
    ]
    for g_head, g_body in glossary_items:
        pdf.set_font("Arial", "B", 10)
        pdf.cell(0, 7, safe_str(g_head), ln=True)
        pdf.set_font("Arial", "", 9)
        pdf.multi_cell(0, 5, safe_str(g_body))
        pdf.ln(2)

    pdf_out = pdf.output(dest='S')
    return pdf_out.encode('latin-1') if isinstance(pdf_out, str) else bytes(pdf_out)

# --- THE VALUATION ENGINE ---
def estimate_value(row):
    name = str(row['Vehicle Name']).lower()
    vehicle_type = str(row['Type']).lower()
    MODEL_PRICES = {
        'f-150': 60000, 'f 150': 60000, 'f150': 60000, 'super duty': 80000,
        'ranger': 40000, 'maverick': 28000, 'bronco': 50000, 'explorer': 48000,
        'expedition': 70000, 'edge': 40000, 'escape': 32000, 'mustang': 45000,
        'mach-e': 55000, 'silverado': 55000, 'sierra': 62000, 'colorado': 38000,
        'tahoe': 72000, 'suburban': 78000, 'yukon': 78000, 'escalade': 110000,
        'corvette': 90000, 'camaro': 42000, 'blazer': 40000, 'equinox': 32000,
        'traverse': 42000, 'malibu': 28000, 'grand cherokee': 55000, 'wagoneer': 80000,
        'wrangler': 48000, 'gladiator': 50000, 'ram 1500': 60000, 'durango': 50000,
        'charger': 40000, 'challenger': 45000, 'pacifica': 45000, 'compass': 30000,
        'tundra': 58000, 'tacoma': 42000, '4runner': 50000, 'sequoia': 75000,
        'highlander': 45000, 'rav4': 35000, 'camry': 30000, 'corolla': 25000,
        'prius': 32000, 'sienna': 48000, 'pilot': 48000, 'passport': 42000,
        'cr-v': 36000, 'odyssey': 45000, 'civic': 28000, 'accord': 32000,
        '3 series': 50000, '5 series': 65000, 'x3': 55000, 'x5': 75000, 'x7': 90000,
        'c-class': 52000, 'e-class': 70000, 'gle': 75000, 'gls': 95000,
        'rx': 58000, 'nx': 48000, 'es': 48000, 'gx': 70000, 'lx': 100000,
        'gv80': 70000, 'gv70': 55000, 'telluride': 48000, 'palisade': 50000,
        'lyriq': 65000, 'optiq': 54000, 'ct4': 40000, 'ct5': 50000, 'xt4': 40000, 'xt5': 50000
    }
    BRAND_DEFAULTS = {
        'cadillac': 65000, 'mercedes': 70000, 'bmw': 65000, 'audi': 60000,
        'lexus': 60000, 'lincoln': 65000, 'land rover': 85000, 'porsche': 100000,
        'ford': 45000, 'chevrolet': 40000, 'gmc': 55000, 'ram': 55000,
        'jeep': 45000, 'dodge': 40000, 'toyota': 38000, 'honda': 35000,
        'nissan': 32000, 'hyundai': 30000, 'kia': 30000, 'subaru': 32000,
        'volkswagen': 32000, 'mazda': 32000, 'volvo': 55000
    }
    baseline = 35000
    found_model = False
    sorted_models = sorted(MODEL_PRICES.keys(), key=len, reverse=True)
    for model in sorted_models:
        if re.search(r'\b' + re.escape(model) + r'\b', name):
            baseline = MODEL_PRICES[model]
            found_model = True
            break
    if not found_model:
        for brand, price in BRAND_DEFAULTS.items():
            if brand in name:
                baseline = price
                break
    if 'new' in vehicle_type:
        return int(baseline)
    
    eastern = pytz.timezone('US/Eastern')
    year_match = re.search(r'\d{4}', name)
    year = int(year_match.group(0)) if year_match else 2025
    current_year = datetime.datetime.now(eastern).year + 1
    age = current_year - year
    if age <= 0:
        value = baseline
    else:
        value = baseline * 0.85
        for _ in range(age - 1):
            value = value * 0.90
    return int(value)

def get_price_tier(price):
    if price < 30000: return "Budget (<$30k)"
    if price < 60000: return "Core ($30k-$60k)"
    return "Premium ($60k+)"

# --- THE SCANNING HELPERS ---
def get_year(url):
    match = re.search(r'(?:^|[^0-9])((?:19|20)\d{2})(?:$|[^0-9])', str(url))
    return match.group(1) if match else None

def extract_vin(url):
    try:
        path = urlparse(str(url)).path.upper().strip('/')
        match = re.search(r'(?:^|[-/])([A-HJ-NPR-Z0-9]{17})(?:[-/\.]|$)', path)
        if match: return match.group(1)
        # Fallback: only accept a real 17-char VIN (excludes I/O/Q). A shorter
        # token is NOT a VIN — grabbing it would cause false "SOLD" results, so
        # we return N/A and let the year-based logic handle it instead.
        blocks = re.findall(r'[A-HJ-NPR-Z0-9]{17}', path)
        if blocks: return blocks[-1]
        return "N/A"
    except:
        return "N/A"

def extract_type(url):
    u = str(url).lower()
    if 'used' in u and 'new' not in u: return 'Used'
    if 'new' in u and 'used' not in u: return 'New'
    if '/used' in u or '-used-' in u or '=used' in u or 'preowned' in u: return 'Used'
    if '/new' in u or '-new-' in u or '=new' in u: return 'New'
    return 'New' if re.search(r'202[5-7]', u) else 'Used'

def smart_dealer_name(url, is_multi=False):
    u = str(url).lower()
    if not u.startswith('http'):
        u = 'http://' + u
    netloc = urlparse(u).netloc
    netloc = re.sub(r'^(www\.|shop\.|inventory\.|cars\.)', '', netloc)
    s = netloc.split('.')[0]
    
    brands = ['honda', 'toyota', 'ford', 'chevrolet', 'chevy', 'nissan', 'jeep', 'chrysler', 'dodge', 'ram',
              'hyundai', 'kia', 'vw', 'volkswagen', 'bmw', 'mercedes', 'audi', 'lexus', 'acura',
              'infiniti', 'subaru', 'mazda', 'volvo', 'porsche', 'buick', 'gmc', 'cadillac', 'lincoln', 'mitsubishi']
    
    modifiers = ['of', 'auto', 'group', 'cars', 'motors', 'dealers', 'cares', 'center', 'city', 'town', 'resale', 'classic', 'one', 'mile', 'superstore']
    
    for mod in modifiers:
        s = s.replace(mod, f" {mod} ")
        
    brands.sort(key=len, reverse=True)
    has_brand = False
    
    for brand in brands:
        if brand in s:
            if brand == 'ram' and ('paramus' in s or 'framingham' in s):
                continue
            if brand == 'ford' and ('oxford' in s or 'crawford' in s or 'stratford' in s):
                continue
            s = s.replace(brand, f" {brand} ")
            has_brand = True
            
    s = re.sub(r'\s+', ' ', s).strip().title()
    s = s.replace(" Vw", " VW").replace(" Bmw", " BMW").replace(" Gmc", " GMC")
    
    if is_multi and not has_brand:
        if any(x in s for x in ["Cares", "Community", "Gives"]):
            s += " (Community)"
        elif any(x in s for x in ["Resale", "Classic", "Used", "Preowned"]):
            pass 
        else:
            s += " (Group/Central Site)"
            
    return s

def clean_name_universal(url):
    year = get_year(url)
    if not year: return "Unknown Vehicle"
    
    parsed = urlparse(str(url))
    path_str = parsed.path + (("?" + parsed.query) if parsed.query else "")
    
    brands = ['Jeep', 'Ford', 'Gmc', 'Toyota', 'Dodge', 'Ram', 'Chrysler', 'Chevrolet', 'Honda', 'Nissan', 'Hyundai', 'Kia', 'Bmw', 'Lexus', 'Volvo', 'Volkswagen', 'Subaru', 'Mazda', 'Mercedes', 'Audi', 'Cadillac', 'Buick', 'Acura', 'Infiniti', 'Lincoln', 'Land Rover', 'Jaguar', 'Porsche', 'Mini']
    make = ""
    for b in brands:
        if b.lower() in path_str.lower():
            make = b.title()
            break
            
    parts = path_str.split(str(year), 1)
    rest = parts[-1] if len(parts) > 1 else path_str
    
    rest = rest.replace('/', ' ').replace('-', ' ').replace('+', ' ').replace('.htm', '').replace('.html', '')
    tokens = rest.split()
    
    junk = ['Baltimore', 'Ephrata', 'Md', 'Maryland', 'Heritage', 'Twin', 'Pine', 'Wholesale', 'New', 'Used', 'Preowned', 'Inventory', 'Parts', 'Service', 'Finance', 'Global', 'Incentives', 'Offers', 'Suv', 'Truck', 'Coupe', 'Sedan', 'Vehicle', 'Vehicles']
    clean_tokens = [t for t in tokens if not (len(t) > 10 and any(c.isdigit() for c in t)) and t.title() not in junk and t.title() != make]
    
    clean_tokens = [t for t in clean_tokens if t != str(year)]
    
    return f"{year} {make} {' '.join(clean_tokens)}".title().strip()

def categorize(u):
    u = str(u).lower()
    
    if any(x in u for x in ['thank', 'confirm', 'success']):
        return 'Online Conversions'
        
    if u.endswith('.com/') or u.endswith('.com'): return 'Homepage'
    if any(x in u for x in ['service', 'parts', 'collision', 'appointment', 'maintenance']): return 'Service'
    
    if any(x in u for x in ['incentive', 'offers', 'specials', 'promotions', 'rebate']): 
        return 'Incentives/Offers'
        
    if 'index.htm' in u or u.endswith('searchnew.aspx') or u.endswith('searchused.aspx') or u.endswith('searchall.aspx'):
        if 'new' in u: return 'New Car Search'
        if 'used' in u or 'preowned' in u: return 'Used Car Search'
        return 'General Search'
        
    if get_year(u): return 'VDP'
    
    if any(x in u for x in ['search', 'inventory', 'vehicles']):
        if 'new' in u: return 'New Car Search'
        if 'used' in u or 'preowned' in u: return 'Used Car Search'
        return 'General Search'
        
    return 'Other'

from collections import defaultdict

# ======================================================================
# DEALER INSPIRE BATCHED VIN RESOLVER
# Checking availability one Algolia query per VIN means thousands of HTTP calls
# (which ground the Basil group report to a halt). Algolia's multi-query endpoint
# lets us pack ~50 VIN lookups into ONE request, and each lookup only needs the
# match count (hitsPerPage=0, no car data), so the payload is tiny. This turns
# ~10,000 requests into a couple hundred and works for a store of any size —
# no 1,000-record pagination ceiling. Runs on algolia.net, never the firewall.
# ======================================================================

def resolve_algolia_vins(domain, cfg, vins, session):
    """
    Resolve a list of VINs for one store via batched Algolia multi-queries.
    Returns {vin: "Available" | "SOLD (Not in Dealer Database)"} for whatever
    resolved; VINs left out (on failure) fall back to the per-VIN path in scan_url.
    """
    app_id, api_key, index_name = cfg['app_id'], cfg['api_key'], cfg['index']
    endpoint = f"https://{app_id.lower()}-dsn.algolia.net/1/indexes/*/queries"
    headers = {
        "x-algolia-application-id": app_id,
        "x-algolia-api-key": api_key,
        "Content-Type": "application/json",
        "Referer": f"https://www.{domain}/",
        "Origin": f"https://www.{domain}",
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                      "(KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36",
    }
    out = {}
    vins = list(vins)
    BATCH = 50
    try:
        for i in range(0, len(vins), BATCH):
            chunk = vins[i:i + BATCH]
            body = {"requests": [
                {"indexName": index_name, "params": f"query={v}&hitsPerPage=0"}
                for v in chunk
            ]}
            r = session.post(endpoint, headers=headers, json=body, timeout=15)
            if r.status_code != 200:
                break  # leave the rest unresolved -> scan_url per-VIN fallback
            results = r.json().get("results", [])
            for v, res in zip(chunk, results):
                nb = res.get("nbHits", 0)
                out[v] = "Available" if nb > 0 else "SOLD (Not in Dealer Database)"
    except Exception:
        pass
    return out

# ======================================================================
# MARKETCHECK EXTERNAL FALLBACK
# For dealers that are locked (API rejected + firewall-blocked) so neither our
# own API nor the scraper can reach them, MarketCheck's daily-crawled inventory
# answers "is this VIN still listed?". Only fires for cars the normal scan left
# as errors, and only if a key is set in Streamlit secrets.
# ======================================================================

def _find_secret(obj, name, depth=0):
    """Recursively search a secrets mapping for a key by name."""
    if depth > 4:
        return ""
    try:
        v = obj.get(name) if hasattr(obj, "get") else None
        if v:
            return str(v).strip()
    except Exception:
        pass
    try:
        for k in obj:
            sub = obj[k]
            if hasattr(sub, "get") or hasattr(sub, "__getitem__"):
                found = _find_secret(sub, name, depth + 1)
                if found:
                    return found
    except Exception:
        pass
    return ""

def get_marketcheck_key():
    """
    Read the MarketCheck key from Streamlit secrets. Prefers the top level (the
    correct place) but searches nested sections too, so a key mistakenly placed
    under a header like [connections.gsheets] is still found. '' if missing.
    """
    try:
        return _find_secret(st.secrets, "MARKETCHECK_KEY")
    except Exception:
        return ""

# Sold-detection endpoint. CONFIRMED BY MARKETCHECK (Jul 2026):
#  - /v2/dealerships/inventory = "Dealer Inventory Syndication API": $1.00/call
#    on PAID tiers, but NO DATA FEES on the FREE tier (10 listings/page there).
#  - /v2/search/car/active = "Inventory Search API" ($0.002/call): officially
#    meant for zip/lat-long searches; source= filtering works in isolation but
#    is unsupported and flaky mid-run (per-VIN vin= fallback covers it).
# FREE-TIER CONFIG (the $0/month setup): set in settings —
#   MARKETCHECK_ENDPOINT = "https://api.marketcheck.com/v2/dealerships/inventory"
#   MARKETCHECK_PAGE_SIZE = 10, MARKETCHECK_PAGE_CAP = 500,
#   MARKETCHECK_BUDGET = 75, MARKETCHECK_MONTHLY_CAP = 450
# On PAID tiers keep the default below (search endpoint, pennies per run).
_MC_DEFAULT_ENDPOINT = "https://api.marketcheck.com/v2/search/car/active"
def _mc_endpoint():
    try:
        v = str(_find_secret(st.secrets, "MARKETCHECK_ENDPOINT")).strip()
        return v if v.startswith("http") else _MC_DEFAULT_ENDPOINT
    except Exception:
        return _MC_DEFAULT_ENDPOINT
MC_ENDPOINT = _mc_endpoint()

def _mc_setting_int(name, default):
    """Read an integer app setting from secrets (any nesting); default if absent."""
    try:
        v = _find_secret(st.secrets, name)
        return int(str(v).strip()) if str(v).strip() else default
    except Exception:
        return default

# Pagination ceiling per inventory slice. Free plan = 500 rows. After upgrading
# to Basic, set MARKETCHECK_PAGE_CAP = 1500 in settings and every pull gets
# cheaper (most dealers fit in one un-sliced pass).
MC_PAGE_CAP = _mc_setting_int("MARKETCHECK_PAGE_CAP", 500)

# Hard spending ceiling per report. A single run can never use more market
# lookups than this, no matter what. Override with MARKETCHECK_BUDGET in
# settings (suggested: 75 on the free plan, 150 on Basic).
MC_BUDGET = _mc_setting_int("MARKETCHECK_BUDGET", 75)

# Optional month-wide ceiling across ALL reports (0 = disabled). For the free
# plan set MARKETCHECK_MONTHLY_CAP = 450 in settings: runs stop spending when
# the month's total (from the UsageStats log) reaches the cap, keeping a
# reserve under the plan's 500-call allowance.
MC_MONTH_CAP = _mc_setting_int("MARKETCHECK_MONTHLY_CAP", 0)

# Dealer-inventory memory: complete MarketCheck store pulls are remembered in
# the Google Sheet ("MCCache" tab) and reused for this many days, so repeat
# reports on the same blocked dealer cost ZERO lookups. 0 disables.
MC_CACHE_DAYS = _mc_setting_int("MARKETCHECK_CACHE_DAYS", 7)

# Cars fetched per lookup. Free plan serves small pages regardless of what's
# requested, so the safe default is 50. Basic honors larger pages — set
# MARKETCHECK_PAGE_SIZE = 200 in settings and big dealers cost ~1/4 the calls.
MC_PAGE_SIZE = _mc_setting_int("MARKETCHECK_PAGE_SIZE", 50)

class _MCStop(Exception):
    """Raised by the governor to halt the market-lookup phase immediately."""
    def __init__(self, reason):
        super().__init__(reason)
        self.reason = reason   # "budget" | "limit" | "errors"

class _MCGovernor:
    """
    Gatekeeper for every market lookup. Enforces the per-run budget, paces
    requests under the plan's rate limit, and halts the phase instantly on
    quota/auth refusals or repeated failures — so a run can never hang and can
    never drain the monthly allowance.
    """
    def __init__(self, session, budget, month_used=0, month_cap=0):
        self._s = session
        self.budget = budget
        self.month_used = int(month_used or 0)
        self.month_cap = int(month_cap or 0)
        self.count = 0
        self.stopped = None        # None, or the _MCStop reason
        self.last_error = None     # human-readable detail of the last failure
        self._fail_streak = 0

    def get(self, *a, **k):
        if self.stopped:
            raise _MCStop(self.stopped)
        if self.count >= self.budget:
            self.stopped = "budget"
            raise _MCStop("budget")
        if self.month_cap and (self.month_used + self.count) >= self.month_cap:
            self.stopped = "monthly"
            raise _MCStop("monthly")
        time.sleep(0.22)           # pace safely under 5 requests/second
        self.count += 1
        try:
            r = self._s.get(*a, **k)
        except _MCStop:
            raise
        except Exception as e:
            self.last_error = f"{type(e).__name__}: {str(e)[:160]}"
            self._fail_streak += 1
            if self._fail_streak >= 3:
                self.stopped = "errors"
                raise _MCStop("errors")
            raise
        if r.status_code in (401, 403, 429):
            # Key rejected / quota exhausted / rate refusal: stop the whole
            # phase now rather than burning calls into a wall.
            self.stopped = "limit"
            raise _MCStop("limit")
        if r.status_code != 200:
            try:
                body = str(r.text)[:200]
            except Exception:
                body = ""
            self.last_error = f"HTTP {r.status_code}: {body}"
            self._fail_streak += 1
            if self._fail_streak >= 3:
                self.stopped = "errors"
                raise _MCStop("errors")
        else:
            self._fail_streak = 0
        return r

def _mc_cache_load(conn):
    """Read the MCCache tab -> {domain: {active, absent, complete, age}}. Best-effort."""
    out = {}
    try:
        df = conn.read(worksheet="MCCache", ttl=0)
        now = datetime.datetime.now(datetime.timezone.utc)
        for _, row in df.iterrows():
            dom = str(row.get("Domain", "")).strip().lower()
            try:
                ts = pd.to_datetime(str(row.get("Timestamp", "")), utc=True)
                age = (now - ts.to_pydatetime()).total_seconds() / 86400.0
            except Exception:
                continue
            def _vset(cell):
                return {v for v in (x.strip().upper() for x in str(cell or "").split(","))
                        if len(v) == 17}
            active = _vset(row.get("ActiveVINs", ""))
            absent = _vset(row.get("AbsentVINs", ""))
            complete = str(row.get("Complete", "")).strip().upper() in ("TRUE", "1", "YES")
            if dom and (active or absent) and age >= 0:
                out[dom] = {"active": active, "absent": absent,
                            "complete": complete, "age": age}
    except Exception:
        pass
    return out

def _mc_cache_save(conn, updates, max_age_days=7):
    """
    Upsert {domain: {active, absent, complete}} into MCCache. A partial (per-VIN)
    result merges into a fresh existing row instead of replacing it; a complete
    store snapshot replaces the row and restarts its clock. Never raises.
    """
    if not updates:
        return
    try:
        cols = ["Domain", "ActiveVINs", "AbsentVINs", "Complete", "Count", "Timestamp"]
        try:
            df = conn.read(worksheet="MCCache", ttl=0)
            if df is None or not isinstance(df, pd.DataFrame):
                df = pd.DataFrame(columns=cols)
        except Exception:
            df = pd.DataFrame(columns=cols)
        existing = _mc_cache_load(conn) if len(df) else {}
        rows = {}
        for _, r in df.iterrows():
            d = str(r.get("Domain", "")).strip().lower()
            if d:
                rows[d] = {c: r.get(c, "") for c in cols}
        now = str(datetime.datetime.now(datetime.timezone.utc))
        for dom, upd in updates.items():
            active = set(upd.get("active") or set())
            absent = set(upd.get("absent") or set())
            complete = bool(upd.get("complete"))
            prev = existing.get(dom)
            ts = now
            if prev and prev["age"] <= max_age_days and not complete:
                # Merge partial knowledge into the fresh existing row; keep its
                # clock so staleness is anchored to the oldest information.
                active |= prev["active"]
                absent |= prev["absent"] - active
                complete = prev["complete"]
                old_row = rows.get(dom, {})
                ts = old_row.get("Timestamp", now) or now
            if complete:
                absent = set()   # a complete snapshot implies absence for the rest
            if len(active) + len(absent) > 2500:   # keep cells under Sheets limits
                continue
            rows[dom] = {"Domain": dom,
                         "ActiveVINs": ",".join(sorted(active)),
                         "AbsentVINs": ",".join(sorted(absent)),
                         "Complete": "TRUE" if complete else "FALSE",
                         "Count": len(active) + len(absent),
                         "Timestamp": ts}
        data = pd.DataFrame(list(rows.values()), columns=cols)
        try:
            conn.update(worksheet="MCCache", data=data)
        except Exception:
            conn.create(worksheet="MCCache", data=data)   # tab doesn't exist yet
    except Exception:
        pass

def _mc_request(domain, api_key, session, year=None, year_range=None, car_type=None, price_range=None, rows=50, start=0):
    """One MarketCheck dealer-inventory request. Returns (num_found, listings) or (None, None)."""
    on_search = (MC_ENDPOINT == _MC_DEFAULT_ENDPOINT)
    if on_search:
        # The Inventory Search API hard-caps rows at 50 and 422s on more —
        # clamp here so a MARKETCHECK_PAGE_SIZE=200 setting (fine on the old
        # syndication endpoint) can't poison every request.
        try:
            rows = min(int(rows), 50)
        except Exception:
            rows = 50
    params = {"api_key": api_key, "source": domain, "rows": str(rows), "start": str(start)}
    if year:
        params["year"] = str(year)
    if year_range:
        translated = False
        if on_search:
            # The search endpoint has no year_range param; send a comma list
            # of years instead (equivalent filter, universally supported).
            try:
                lo, hi = (int(x) for x in str(year_range).split("-"))
                if 0 <= hi - lo <= 30:
                    params["year"] = ",".join(str(y) for y in range(lo, hi + 1))
                    translated = True
            except Exception:
                pass
        if not translated:
            params["year_range"] = year_range
    if car_type:
        params["car_type"] = car_type
    if price_range:
        params["price_range"] = price_range
    try:
        r = session.get(MC_ENDPOINT, params=params, timeout=25)
        if r.status_code != 200:
            return (None, None)
        data = r.json()
        if "num_found" not in data:
            if hasattr(session, "last_error"):
                session.last_error = f"200 OK but unexpected body: {str(data)[:140]}"
            return (None, None)
        return (data.get("num_found", 0), data.get("listings", []))
    except _MCStop:
        raise   # governor halt: must reach the orchestrator, never swallowed
    except Exception as e:
        if hasattr(session, "last_error"):
            session.last_error = f"response error: {type(e).__name__}: {str(e)[:140]}"
        return (None, None)

def _mc_pull(domain, api_key, session, year=None, year_range=None, car_type=None, price_range=None):
    """
    Pull a slice's VINs, paging in 50s. Aborts early (no paging) if the slice is
    over the free-tier cap, so an over-cap slice costs ONE call, not ten.
    Returns (vin_set, complete, num_found): complete means we got the whole slice
    (so an absent VIN is genuinely sold); num_found is the slice's total.
    """
    nf, listings = _mc_request(domain, api_key, session, year=year, year_range=year_range,
                               car_type=car_type, price_range=price_range, rows=MC_PAGE_SIZE, start=0)
    if nf is None:
        return (set(), False, None)
    vins = set()
    for lst in listings:
        v = str(lst.get("vin", "")).upper().strip()
        if v:
            vins.add(v)
    if nf > MC_PAGE_CAP:
        return (vins, False, nf)   # over cap -> caller sub-slices; don't waste pages
    start = len(listings)
    while start < nf and start < MC_PAGE_CAP:
        nf2, more = _mc_request(domain, api_key, session, year=year, year_range=year_range,
                                car_type=car_type, price_range=price_range, rows=MC_PAGE_SIZE, start=start)
        if nf2 is None or not more:
            break
        for lst in more:
            v = str(lst.get("vin", "")).upper().strip()
            if v:
                vins.add(v)
        start += len(more)
    return (vins, start >= nf, nf)

def _mc_apply(out, items, vin_set, complete):
    """Assign Available/Sold to report items from a pulled slice."""
    for url, vin in items:
        if vin.upper().strip() in vin_set:
            out[url] = "Available"
        elif complete:
            out[url] = "SOLD (Not in Market Inventory)"
        # partial slice -> leave the car's original result (no false sold)

def _mc_pull_priced(domain, api_key, session, year, car_type, lo, hi, depth=0, parent_nf=None):
    """
    Pull a car_type slice within a price band [lo, hi]. If that band is STILL over
    the cap, split it at the midpoint and recurse — so any distribution eventually
    fits. Stops splitting if a narrower band doesn't actually shrink the results
    (filter ignored by the data), so calls are never wasted on useless recursion.
    Returns (vin_set, complete).
    """
    vins, complete, nf = _mc_pull(domain, api_key, session, year=year,
                                  car_type=car_type, price_range=f"{lo}-{hi}")
    if nf is None:
        return (set(), False)
    if complete:
        return (vins, True)
    if parent_nf is not None and nf >= parent_nf:
        return (vins, False)   # band didn't narrow anything -> splitting is futile
    if depth >= 4 or (hi - lo) <= 1000:
        return (vins, False)   # can't split further -> partial (never false-sold)
    mid = (lo + hi) // 2
    v1, c1 = _mc_pull_priced(domain, api_key, session, year, car_type, lo, mid, depth + 1, nf)
    v2, c2 = _mc_pull_priced(domain, api_key, session, year, car_type, mid, hi, depth + 1, nf)
    return (vins | v1 | v2, c1 and c2)

# Starting price bands for a car_type slice that's over the cap (tuned for autos;
# any band still too large is split further by _mc_pull_priced).
MC_PRICE_BANDS = [(0, 45000), (45000, 60000), (60000, 80000), (80000, 100000000)]

def _mc_resolve_year(domain, api_key, session, year, yr_items, out):
    """Resolve one model year; if over the cap, split by car_type, then by price."""
    vins, complete, nf = _mc_pull(domain, api_key, session, year=year)
    if nf is None:
        return
    if complete:
        _mc_apply(out, yr_items, vins, True)
        return
    # Over cap -> split by car_type (each piece usually fits under the cap).
    merged = set(vins)
    all_complete = True
    for ct in ("new", "used", "certified"):
        vs, cp, nf2 = _mc_pull(domain, api_key, session, year=year, car_type=ct)
        if nf2 is None:
            all_complete = False
            continue
        if cp:
            merged |= vs
        else:
            # A single car_type still over the cap (e.g. 694 new units) -> split
            # it by price band, recursing on any band that's itself too large.
            for lo, hi in MC_PRICE_BANDS:
                pvs, pcp = _mc_pull_priced(domain, api_key, session, year, ct, lo, hi,
                                           parent_nf=nf2)
                merged |= pvs
                all_complete = all_complete and pcp
    _mc_apply(out, yr_items, merged, all_complete)

def _mc_band(domain, api_key, session, band_items, out):
    """
    Resolve a group of older report years as ONE year_range band (cheap when the
    years are sparse). If the band is over the cap, fall back to per-year.
    band_items = [(url, vin, year_int)].
    """
    if not band_items:
        return
    years = [y for _, _, y in band_items]
    lo, hi = min(years), max(years)
    vins, complete, nf = _mc_pull(domain, api_key, session, year_range=f"{lo}-{hi}")
    if nf is None:
        return
    if complete:
        _mc_apply(out, [(u, v) for u, v, _ in band_items], vins, True)
        return
    # Band over the cap -> resolve each year on its own.
    per_year = defaultdict(list)
    for u, v, y in band_items:
        per_year[str(y)].append((u, v))
    for yr, its in per_year.items():
        _mc_resolve_year(domain, api_key, session, yr, its, out)

def _mc_vin_batch(chunk, api_key, session):
    """
    One vins= batch lookup for [(url, vin)]. Returns ({url: status}, error).
    error is None on success; a string when the filter isn't honored or the
    call fails (caller decides whether to fall back).
    """
    params = {"api_key": api_key, "vins": ",".join(v for _, v in chunk),
              "rows": "50", "start": "0"}
    r = session.get(MC_ENDPOINT, params=params, timeout=25)
    if r.status_code != 200:
        try:
            body = str(r.text)[:160]
        except Exception:
            body = ""
        return ({}, f"vins= batch: HTTP {r.status_code}: {body}")
    data = r.json()
    nf = int(data.get("num_found", 0) or 0)
    listings = data.get("listings", [])
    # Same VIN can be listed by a few sources, but num_found beyond ~20x the
    # batch means the vins filter was ignored (national results). Refuse to
    # interpret that — never risk a false Sold.
    if nf > len(chunk) * 20:
        return ({}, f"vins= filter not honored (num_found={nf} for {len(chunk)} VINs)")
    active = set()
    for lst in listings:
        vv = str(lst.get("vin", "")).upper().strip()
        if vv:
            active.add(vv)
    start = len(listings)
    while start < nf and start < 400:
        r2 = session.get(MC_ENDPOINT, params={**params, "start": str(start)}, timeout=25)
        if r2.status_code != 200:
            break
        more = r2.json().get("listings", [])
        if not more:
            break
        for lst in more:
            vv = str(lst.get("vin", "")).upper().strip()
            if vv:
                active.add(vv)
        start += len(more)
    return ({u: ("Available" if v in active else "SOLD (Not in Market Inventory)")
             for u, v in chunk}, None)

def _mc_vin_single(url, vin, api_key, session):
    """
    One vin= lookup. Returns (status_or_None, error). num_found for a single
    VIN must be tiny; a huge count means the filter is ignored -> error.
    """
    r = session.get(MC_ENDPOINT, params={"api_key": api_key, "vin": vin,
                                         "rows": "10", "start": "0"}, timeout=25)
    if r.status_code != 200:
        try:
            body = str(r.text)[:160]
        except Exception:
            body = ""
        return (None, f"vin= single: HTTP {r.status_code}: {body}")
    data = r.json()
    nf = int(data.get("num_found", 0) or 0)
    if nf > 50:
        return (None, f"vin= filter not honored (num_found={nf} for 1 VIN)")
    active = any(str(l.get("vin", "")).upper().strip() == vin
                 for l in data.get("listings", []))
    return ("Available" if active else "SOLD (Not in Market Inventory)", None)

def _mc_resolve_vins(domain, items, api_key, session, cache_out=None):
    """
    VIN-based resolution for the search endpoint: ask the market about the
    exact VINs in the report instead of pulling the dealer's whole store.
    Adaptive: tries a vins= batch first (~40 VINs per call); if that filter
    isn't honored by this account, falls back to per-VIN vin= lookups (one
    call each). If neither filter works, every failure detail is recorded
    and the dealer is left unresolved — never guessed.
    items = [(url, vin, year)]. Returns {url: status}.
    """
    out = {}
    valid = []
    for u, v, _y in items:
        vv = str(v or "").upper().strip()
        if len(vv) == 17:
            valid.append((u, vv))
    if not valid:
        return out
    CHUNK = 40
    # vins= is disabled: the 7/20 diagnostic returned num_found=211 for a
    # SINGLE VIN (wrong semantics — likely matching historical listings),
    # while vin= returned exactly 1. Per-VIN is production-proven; the batch
    # machinery is kept below in case MarketCheck fixes the param.
    mode = "single"
    batch_err = None
    for i in range(0, len(valid), CHUNK):
        chunk = valid[i:i + CHUNK]
        if mode in (None, "batch"):
            try:
                res, err = _mc_vin_batch(chunk, api_key, session)
            except _MCStop:
                raise
            except Exception as e:
                res, err = {}, f"vins= batch: {type(e).__name__}"
            if err is None:
                mode = "batch"
                out.update(res)
                continue
            if mode == "batch":
                # Batch worked earlier but failed now: transient — skip chunk.
                if hasattr(session, "last_error"):
                    session.last_error = err
                continue
            batch_err = err   # first chunk failed -> try per-VIN fallback
            mode = "single"
        # Per-VIN mode (fallback)
        for u, v in chunk:
            try:
                status, err = _mc_vin_single(u, v, api_key, session)
            except _MCStop:
                raise
            except Exception as e:
                status, err = None, f"vin= single: {type(e).__name__}"
            if err is not None:
                if hasattr(session, "last_error"):
                    session.last_error = (err if not batch_err
                                          else f"{batch_err} | then {err}")
                if "not honored" in err or "HTTP 4" in err:
                    break   # systemic: stop burning calls on this dealer
                continue
            out[u] = status
    # Remember what the per-VIN checks learned (partial dealer knowledge).
    if cache_out is not None and out:
        entry = cache_out.setdefault(domain, {"active": set(), "absent": set(),
                                              "complete": False})
        by_url = {u: v for u, v in valid}
        for u, status in out.items():
            v = by_url.get(u)
            if not v:
                continue
            if status == "Available":
                entry["active"].add(v)
            elif str(status).startswith("SOLD"):
                entry["absent"].add(v)
    return out

def resolve_dealer_marketcheck(domain, items, api_key, session, cache_out=None):
    """
    Resolve one locked dealer's unresolved cars against MarketCheck, minimizing
    calls. items = [(url, vin, year)]. Each slice is size-checked as it's pulled
    (aborting after one call if it's over the cap) so we never waste pages:
      • Small dealers (under the cap): the whole store in one pass.
      • Big dealers, tiered by age: newest years individually (dense; split by
        car_type if a single year is over the cap), older years bundled into bands.
    Returns {url: "Available" | "SOLD (Not in Market Inventory)"}.
    """
    out = {}
    # On the search endpoint: store-pull first (the 7/20 diagnostic confirmed
    # source= is honored there — Hamby returned num_found=75, dealer-sized),
    # which resolves a whole dealer in ceil(inventory/50) calls. The VIN path
    # (proven in production 7/20: exact match to baseline) is the fallback for
    # errors, unfiltered-looking responses, and whatever a partial pull left.
    if MC_ENDPOINT == _MC_DEFAULT_ENDPOINT:
        vins_s, complete_s, nf_s = _mc_pull(domain, api_key, session)
        if nf_s is not None and 0 <= nf_s <= 25000:
            if complete_s:
                if cache_out is not None:
                    cache_out[domain] = {"active": set(vins_s), "absent": set(),
                                         "complete": True}
                _mc_apply(out, [(u, v) for u, v, _ in items], vins_s, True)
                return out
            # Partial pull (over the page cap): found VINs are safely
            # Available; resolve the remainder per-VIN.
            _mc_apply(out, [(u, v) for u, v, _ in items], vins_s, False)
            if cache_out is not None:
                cache_out[domain] = {"active": set(vins_s), "absent": set(),
                                     "complete": False}
            remaining = [(u, v, y) for u, v, y in items if u not in out]
            out.update(_mc_resolve_vins(domain, remaining, api_key, session,
                                        cache_out=cache_out))
            return out
        if nf_s is not None and nf_s > 25000 and hasattr(session, "last_error"):
            session.last_error = (f"source filter returned num_found={nf_s} for "
                                  f"{domain} — falling back to VIN lookups")
        out.update(_mc_resolve_vins(domain, items, api_key, session,
                                    cache_out=cache_out))
        return out
    # Try the whole store; if it's over the cap this aborts after one call.
    vins, complete, nf = _mc_pull(domain, api_key, session)
    if nf is None:
        return out
    if nf > 25000:
        # No dealership has 25k cars: the source filter isn't filtering.
        # Stop this dealer at ONE call instead of sub-slicing into a storm.
        if hasattr(session, "last_error"):
            session.last_error = (f"source filter not honored for {domain} "
                                  f"(num_found={nf}) — dealer left unresolved")
        return out
    if complete:
        if cache_out is not None:
            cache_out[domain] = {"active": set(vins), "absent": set(),
                                 "complete": True}
        _mc_apply(out, [(u, v) for u, v, _ in items], vins, True)
        return out

    # Big dealer: tier by vehicle age.
    try:
        current_year = datetime.datetime.now().year
    except Exception:
        current_year = max((int(y) for _, _, y in items if str(y).isdigit()), default=2025)
    recent_cut = current_year - 2   # newest 3 model years -> individual
    mid_cut = current_year - 9       # 4-9 years old -> mid band; older -> old band

    recent = defaultdict(list)   # "year" -> [(url, vin)]
    mid_items = []               # [(url, vin, year_int)]
    old_items = []
    for url, vin, year in items:
        if not (year and str(year).isdigit()):
            continue   # no reliable year on a >cap dealer -> leave unresolved (rare)
        y = int(year)
        if y >= recent_cut:
            recent[str(y)].append((url, vin))
        elif y >= mid_cut:
            mid_items.append((url, vin, y))
        else:
            old_items.append((url, vin, y))

    for year, yr_items in recent.items():
        _mc_resolve_year(domain, api_key, session, year, yr_items, out)
    _mc_band(domain, api_key, session, mid_items, out)
    _mc_band(domain, api_key, session, old_items, out)
    return out

def build_all_vin_status(vdp_urls, session):
    """
    Pre-pass: group the report's VINs by vaulted Dealer Inspire domain and resolve
    each store's VINs in batches (concurrently). Returns {domain: {vin: status}}.
    """
    by_domain = defaultdict(set)
    for u in vdp_urls:
        d = urlparse(str(u)).netloc.replace('www.', '').lower()
        cfg = DEALER_API_VAULT.get(d)
        if isinstance(cfg, dict) and not cfg.get('ignore') \
           and cfg.get('app_id') and cfg.get('api_key') and cfg.get('index'):
            v = extract_vin(u)
            if v != "N/A":
                by_domain[d].add(v)
    status = {}
    if not by_domain:
        return status
    with concurrent.futures.ThreadPoolExecutor(max_workers=min(8, len(by_domain))) as ex:
        futs = {ex.submit(resolve_algolia_vins, d, DEALER_API_VAULT[d], vins, session): d
                for d, vins in by_domain.items()}
        for fut in concurrent.futures.as_completed(futs):
            d = futs[fut]
            try:
                status[d] = fut.result()
            except Exception:
                status[d] = {}
    return status

def scan_url(url, session, ignored_domains, ignore_lock, vin_status=None):
    """
    Decide a VDP's status. Order of operations:
      1. HARD IGNORE: if the sheet flags the domain, or we've already been blocked
         by it this run, go straight to the visual scraper. No API call is built.
      2. If there are no usable credentials, use the scraper.
      3. Otherwise fire the Algolia API. On 200 -> read hits. On 401/403/429 (or any
         network error) -> blacklist the domain for the rest of this run and fall
         back to the scraper immediately. We never return a raw API ERROR for a
         blocked API — the scraper always gets a chance.
    """
    domain = urlparse(str(url)).netloc.replace('www.', '').lower()

    # --- 1. NON-DEALER-INSPIRE DEALERS -> SCRAPER ------------------------
    cfg = DEALER_API_VAULT.get(domain)
    sheet_says_ignore = isinstance(cfg, dict) and cfg.get('ignore') is True
    has_creds = (
        isinstance(cfg, dict)
        and cfg.get('app_id') and cfg.get('api_key') and cfg.get('index')
    )
    # Dealers marked ignore, or with no API creds, run on a scrapable platform.
    if sheet_says_ignore or not has_creds:
        return check_universal_status(url, session)

    # --- 2. DEALER INSPIRE (Algolia) DEALERS -----------------------------
    # These sites are firewall-protected, so scraping them is doomed — if the
    # API can't answer, we fail fast rather than dragging the report through a
    # blocked scrape.
    with ignore_lock:
        if domain in ignored_domains:          # this store's key already rejected
            return "ERROR (Inventory Unavailable)"

    vin = extract_vin(url)
    if vin == "N/A":
        return "ERROR (Inventory Unavailable)"

    # Resolved in bulk by the pre-pass?
    pre = (vin_status or {}).get(domain)
    if pre and vin in pre:
        return pre[vin]

    # --- 3. PER-VIN FALLBACK (only if the bulk pass missed this VIN) ------
    app_id, api_key, index_name = cfg['app_id'], cfg['api_key'], cfg['index']
    api_endpoint = f"https://{app_id.lower()}-dsn.algolia.net/1/indexes/{index_name}/query"
    api_headers = {
        "x-algolia-application-id": app_id,
        "x-algolia-api-key": api_key,
        "Content-Type": "application/json",
        "Referer": f"https://www.{domain}/",
        "Origin": f"https://www.{domain}",
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                      "(KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36",
    }
    try:
        resp = session.post(api_endpoint, headers=api_headers,
                            json={"params": f"query={vin}"}, timeout=5)
    except Exception:
        return "ERROR (Inventory Unavailable)"

    if resp.status_code == 200:
        try:
            hits = resp.json().get("nbHits", 0)
        except Exception:
            return "ERROR (Inventory Unavailable)"
        return "Available" if hits > 0 else "SOLD (Not in Dealer Database)"

    # Key rejected / rate-limited -> blacklist this store so the rest fail fast.
    if resp.status_code in (401, 403, 429):
        with ignore_lock:
            ignored_domains.add(domain)
    return "ERROR (Inventory Unavailable)"

# --- THE UNIVERSAL VISUAL SCRAPER ---------------------------------------
def check_universal_status(url, session):
    url = str(url).strip() # STRIP INVISIBLE SPACES FROM CSV
    year = get_year(url)
    vin = extract_vin(url)
    if not year: return "N/A"

    # WAF EVASION HEADERS
    user_agents = [
        'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36',
        'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36',
        'Mozilla/5.0 (Windows NT 10.0; Win64; x64; rv:123.0) Gecko/20100101 Firefox/123.0'
    ]

    try:
        headers = {
            'User-Agent': random.choice(user_agents),
            'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,*/*;q=0.8',
            'Accept-Language': 'en-US,en;q=0.5',
            'Connection': 'keep-alive',
            'Upgrade-Insecure-Requests': '1',
            'Sec-Fetch-Dest': 'document',
            'Sec-Fetch-Mode': 'navigate',
            'Sec-Fetch-Site': 'none',
            'Sec-Fetch-User': '?1'
        }

        response = session.get(url, headers=headers, timeout=5, allow_redirects=True)

        if response.status_code in [403, 406, 429]:
            return f"ERROR (Website Firewall Blocked Scan: {response.status_code})"
            
        if response.status_code in [404, 410]:
            return "SOLD (404 Error)"
            
        orig_base = url.lower().split('?')[0].rstrip('/').replace('https://', '').replace('http://', '').replace('www.', '')
        final_base = response.url.lower().split('?')[0].rstrip('/').replace('https://', '').replace('http://', '').replace('www.', '')
        
        if orig_base != final_base:
            if vin != "N/A" and vin.lower() not in final_base: return "SOLD (HTTP Redirect)"
            if year not in final_base: return "SOLD (HTTP Redirect)"

        text = response.text 
        text_lower = text.lower()
        soup = BeautifulSoup(text, 'html.parser')
        page_title = soup.title.string.strip().lower() if soup.title else ""
        
        # --- SOFT-SOLD / OVERLAY / JSON SCANNER ---
        soft_sold_phrases = [
            "no longer available",
            "this vehicle is sold",
            "vehicle has been sold",
            "currently out of stock",
            "schema.org/outofstock",     
            "schema.org/soldout",        
            '"inventorystatus":"sold"',  
            '"inventorystatus": "sold"', 
            '"vehiclestatus":"sold"',    
            '"vehiclestatus": "sold"',
            '"isavailable":false',
            '"isavailable": false'
        ]
        if any(phrase in text_lower for phrase in soft_sold_phrases):
            return "SOLD (Out of Stock Overlay)"
        # ----------------------------------------
        
        bot_titles = ['just a moment', 'attention required', 'verify you are human', 'access denied', 'pardon our interruption', 'security check']
        if any(b in page_title for b in bot_titles): return "ERROR (Cloudflare Bot Block)"
            
        meta_refresh = soup.find('meta', attrs={'http-equiv': re.compile(r'^refresh$', re.I)})
        if meta_refresh and meta_refresh.get('content'):
            content = meta_refresh['content']
            match = re.search(r'url=([^"\'>\s]+)', content, re.IGNORECASE)
            if match:
                meta_url = match.group(1).lower()
                if vin != "N/A" and vin.lower() not in meta_url:
                    return "SOLD (Meta Refresh Redirect)"
                    
        js_redirects = re.findall(r'window\.location\.(?:replace|href|assign)\s*=\s*["\']([^"\'>]+)["\']', text, re.IGNORECASE)
        for js_url in js_redirects:
            if vin != "N/A" and vin.lower() not in js_url.lower():
                return "SOLD (JS Redirect)"

        if 'not found' in page_title or '404' in page_title or 'error' in page_title: return "SOLD (Page Not Found)"
            
        search_indicators = ['search', 'results', 'all vehicles', 'inventory']
        if any(x in page_title for x in search_indicators) and year not in page_title: return "SOLD (Soft Redirect)"
            
        return "Available"
        
    except requests.exceptions.Timeout: return "ERROR (Timeout)"
    except requests.exceptions.ConnectionError: return "ERROR (Connection Blocked)"
    except Exception as e: return "Available"


# --- UI DASHBOARD ---
st.title("🚗 Auto-Sales Intelligence Agent")

# Sidebar: File Upload
st.sidebar.markdown("### 📥 New Analysis")
uploaded_file = st.sidebar.file_uploader("Upload Traffic Report (CSV)", type=['csv'])

run_analysis_clicked = False
if uploaded_file is not None:
    run_analysis_clicked = st.sidebar.button("🚀 Run Diagnostic Analysis", type="primary", use_container_width=True)
    st.sidebar.info("Upload a CSV and click Run. If we need help locating the dealer's inventory database, we'll give you simple steps to sync it dynamically!")

# --- ADD A DEALER TO THE VAULT (salesperson-friendly) ---
with st.sidebar.expander("🔐 Add Dealer to Vault", expanded=False):
    st.caption("Only needed if a dealer's inventory won't scan on its own. Pick the easiest option that works.")
    tab_auto, tab_book, tab_manual = st.tabs(["⚡ Auto-Detect", "🔖 Bookmarklet", "✍️ Manual"])

    # ---- Option A: Auto-Detect (easiest; may be blocked by strict firewalls) ----
    with tab_auto:
        st.markdown("<span style='font-size:0.85em;color:#555;'>Type the dealer's website and we'll try to find the keys automatically.</span>", unsafe_allow_html=True)
        ad_domain = st.text_input("Dealer website", placeholder="e.g., pittstoyota.com", key="ad_domain")
        if st.button("⚡ Auto-Detect & Save", use_container_width=True, key="ad_btn"):
            dom = normalize_domain(ad_domain)
            if not dom:
                st.warning("Enter the dealer's website first.")
            else:
                with st.spinner(f"Searching {dom} for inventory keys..."):
                    creds = autodetect_algolia(dom)
                if creds:
                    if upsert_vault_row(dom, creds['app_id'], creds['api_key'], creds['index']):
                        st.success(f"✅ {dom} synced automatically!")
                        st.rerun()
                    else:
                        st.error("Found the keys, but couldn't write to the Google Sheet.")
                else:
                    st.warning("Couldn't auto-detect (this dealer's firewall likely blocked us). "
                               "Use the **🔖 Bookmarklet** tab — it runs in your own browser and works even when this doesn't.")

    # ---- Option B: Bookmarklet (most reliable; runs in the salesperson's browser) ----
    with tab_book:
        st.markdown(
            "<span style='font-size:0.85em;color:#555;'>"
            "<b>One-time setup:</b> drag the button below up to your browser's bookmarks bar."
            "</span>", unsafe_allow_html=True)

        safe_href = html.escape(BOOKMARKLET_JS, quote=True)
        components.html(
            f"""
            <div style="font-family:sans-serif;padding:4px 0;">
              <a href="{safe_href}"
                 style="display:inline-block;background:#D70015;color:#fff;text-decoration:none;
                        font-weight:bold;padding:8px 14px;border-radius:6px;border:1px solid #A30010;
                        cursor:grab;font-size:14px;">
                 🔑 Get Dealer Keys
              </a>
              <div style="font-size:12px;color:#777;margin-top:6px;">☝️ Drag me to your bookmarks bar</div>
            </div>
            """,
            height=70,
        )

        st.markdown(
            "<span style='font-size:0.85em;color:#555;'>"
            "<b>Every time:</b> open the dealer's website → click your <b>Get Dealer Keys</b> bookmark. "
            "If it says the Search Key is missing, type any letter in the dealer's inventory search box "
            "so a search runs, then click the bookmark <b>again</b>. It copies a line → paste it below and Save."
            "</span>", unsafe_allow_html=True)

        with st.popover("Can't drag it? Copy the code instead"):
            st.caption("Create a new bookmark, name it 'Get Dealer Keys', and paste this as the URL/address:")
            st.code(BOOKMARKLET_JS, language=None)

        book_paste = st.text_area("Paste the copied line here", placeholder="fivestarfordga.com|APPID|APIKEY|indexname",
                                  key="book_paste", height=80)
        if st.button("💾 Save to Vault", use_container_width=True, key="book_btn"):
            creds = parse_algolia_credentials(book_paste)
            if creds and creds['domain']:
                if upsert_vault_row(creds['domain'], creds['app_id'], creds['api_key'], creds['index']):
                    st.success(f"✅ {creds['domain']} synced to the vault!")
                    st.rerun()
                else:
                    st.error("Couldn't write to the Google Sheet.")
            else:
                st.error("That didn't look right. Click the bookmark on the dealer's site, then paste what it copies.")

    # ---- Option C: Manual (power users / DevTools method) ----
    with tab_manual:
        st.markdown("<span style='font-size:0.85em;color:#555;'>Paste a full Algolia API URL (or any text containing the keys).</span>", unsafe_allow_html=True)
        qd_domain = st.text_input("Dealer website", placeholder="e.g., pittstoyota.com", key="qd_domain")
        qd_url = st.text_area("Algolia API URL / config", placeholder="Paste here...", key="qd_url", height=100)
        if st.button("💾 Save to Vault", use_container_width=True, key="manual_btn"):
            if qd_domain and qd_url:
                creds = parse_algolia_credentials(qd_url, fallback_domain=qd_domain)
                if creds:
                    dom = creds['domain'] or normalize_domain(qd_domain)
                    if upsert_vault_row(dom, creds['app_id'], creds['api_key'], creds['index']):
                        st.success(f"✅ {dom} synced to the vault!")
                        st.rerun()
                    else:
                        st.error("Couldn't write to the Google Sheet. Check the connection and try again.")
                else:
                    st.error("Couldn't find all three keys (App ID, API Key, Index) in what you pasted.")
            else:
                st.warning("Please fill both fields.")

    # ---- Diagnostic: is a dealer's key actually working? ----
    st.divider()
    st.markdown("<span style='font-size:0.85em;color:#555;'><b>🔬 Test a dealer's keys</b> — check if the API actually works (or why it doesn't).</span>", unsafe_allow_html=True)
    test_domain = st.text_input("Domain to test", placeholder="e.g., fivestarfordga.com", key="test_domain")
    test_paste = st.text_area("(optional) paste a bookmarklet line to test before saving", key="test_paste", height=68,
                              placeholder="Leave blank to test what's already saved in the vault")
    if st.button("🔬 Run Key Test", use_container_width=True, key="test_btn"):
        dom = normalize_domain(test_domain)
        creds = parse_algolia_credentials(test_paste, fallback_domain=dom) if test_paste.strip() else None
        source = "pasted line"
        if not creds:
            cfg = DEALER_API_VAULT.get(dom)
            source = "the vault"
            if not dom:
                st.warning("Enter a domain (or paste a line) to test.")
                cfg = None
            elif cfg is None:
                st.error(f"❌ '{dom}' is NOT in the vault — nothing was saved. Add it via Auto-Detect or the Bookmarklet.")
                cfg = None
            elif isinstance(cfg, dict) and cfg.get('ignore'):
                st.error(f"🚫 '{dom}' is set to IGNORE in the sheet, so the app skips the API entirely. "
                         f"Re-save real keys for it to overwrite the IGNORE flag.")
                cfg = None
            if cfg:
                creds = {'app_id': cfg.get('app_id'), 'api_key': cfg.get('api_key'), 'index': cfg.get('index')}
        if creds:
            with st.spinner("Testing keys against Algolia..."):
                ok, headline, detail = test_algolia_creds(creds['app_id'], creds['api_key'], creds['index'], domain=dom)
            (st.success if ok else st.error)(f"{headline}  \n_{detail}_  \n(tested from {source})")

# Main Execution Logic
if run_analysis_clicked:
    df_raw = pd.read_csv(uploaded_file)
    uploaded_file.seek(0)
    
    url_col = 'Page Url' if 'Page Url' in df_raw.columns else 'Page URL' if 'Page URL' in df_raw.columns else df_raw.columns[0]
    df_raw.rename(columns={url_col: 'Page Url'}, inplace=True)
    
    df_raw = df_raw.dropna(subset=['Page Url'])
    df_raw = df_raw[df_raw['Page Url'].astype(str).str.strip() != '']
    
    df_raw['Category'] = df_raw['Page Url'].apply(categorize)
    
    unique_domains_list = df_raw['Page Url'].apply(lambda x: urlparse(str(x)).netloc.replace('www.', '').lower()).unique()
    unique_domains_list = [d for d in unique_domains_list if d]
    is_multi_dealer = len(unique_domains_list) > 1
    
    df_raw['Dealer'] = df_raw['Page Url'].apply(lambda x: smart_dealer_name(x, is_multi_dealer)) 
    vdp_urls = df_raw[df_raw['Category'] == 'VDP']['Page Url'].tolist()
    
    st.info(f"Scanning {len(vdp_urls):,} vehicles for current availability...")
    progress_bar = st.progress(0)
    
    session = requests.Session()
    retry_strategy = Retry(total=3, backoff_factor=1, status_forcelist=[429, 500, 502, 503, 504])
    adapter = HTTPAdapter(max_retries=retry_strategy, pool_connections=60, pool_maxsize=60)
    session.mount('https://', adapter)
    session.mount('http://', adapter)
    
    vdp_results = {}

    # Per-run blacklist of domains whose API blocked us (403/401/429). Once a
    # domain lands here, every remaining VDP for it skips the API entirely.
    ignored_domains = set()
    ignore_lock = threading.Lock()

    # Resolve vaulted dealers' inventory in bulk before the per-car scan.
    with st.spinner("Checking dealer inventories..."):
        vin_status = build_all_vin_status(vdp_urls, session)

    with concurrent.futures.ThreadPoolExecutor(max_workers=30) as executor:
        future_to_url = {
            executor.submit(scan_url, url, session, ignored_domains, ignore_lock, vin_status): url
            for url in vdp_urls
        }
        total_urls = len(vdp_urls)
        # Update the progress bar at most ~100 times, not once per URL — on a
        # 1,000+ car report, per-URL updates flood the UI and slow the run.
        progress_step = max(1, total_urls // 100)
        for i, future in enumerate(concurrent.futures.as_completed(future_to_url)):
            url = future_to_url[future]
            try:
                vdp_results[url] = future.result()
            except Exception:
                # A single malformed URL must never abort a large report.
                vdp_results[url] = "ERROR (Scan Failed)"
            if (i + 1) % progress_step == 0 or (i + 1) == total_urls:
                progress_bar.progress((i + 1) / total_urls)

    # --- EXTERNAL FALLBACK: resolve anything the normal scan couldn't ------
    # Locked dealers (API rejected + firewall-blocked) leave "ERROR" results.
    # If a MarketCheck key is set, resolve those via MarketCheck's daily inventory.
    mc_calls_used = 0
    mc_key = get_marketcheck_key()
    mc_by_domain = defaultdict(list)   # domain -> [(url, vin)]
    for url, res in vdp_results.items():
        if str(res).startswith("ERROR") and "No VIN" not in str(res):
            v = extract_vin(url)
            if v != "N/A":
                dom = urlparse(str(url)).netloc.replace('www.', '').lower()
                mc_by_domain[dom].append((url, v))

    if mc_by_domain:   # there are locked cars to try to rescue
        if not mc_key:
            st.info("A few dealers couldn't be scanned directly. To resolve them, "
                    "the market-inventory lookup needs to be set up in settings.")
        else:
            total_unres = sum(len(v) for v in mc_by_domain.values())
            st.info(f"{len(mc_by_domain)} dealer(s) with {total_unres} vehicle(s) need the "
                    f"market-inventory lookup — this run will use at most {MC_BUDGET} lookups.")
            gov = _MCGovernor(session, MC_BUDGET,
                              month_used=st.session_state.get('mc_month_usage') or 0,
                              month_cap=MC_MONTH_CAP)
            resolved = 0
            stop_reason = None
            cache_hits = 0
            mc_cache_out = {}
            mc_cache = {}
            if MC_CACHE_DAYS > 0 and gsheets_conn is not None:
                mc_cache = _mc_cache_load(gsheets_conn)
            try:
                # Cheap reachability probe (1 call): a bad key or exhausted plan
                # halts here with a clear reason instead of failing silently.
                gov.get("https://api.marketcheck.com/v2/search/car/active",
                        params={"api_key": mc_key, "source": next(iter(mc_by_domain)), "rows": "1"},
                        timeout=15)
                with st.spinner("Checking remaining vehicles against live market inventory..."):
                    # Smallest dealers first: they're the cheapest to finish, so a
                    # tight budget fully resolves as many dealers as possible
                    # before it runs out (whole dealers beat everyone-half-done).
                    for dom, items in sorted(mc_by_domain.items(), key=lambda kv: len(kv[1])):
                        cached = mc_cache.get(dom)
                        if cached is not None and cached["age"] <= MC_CACHE_DAYS:
                            # Serve everything the memory knows at zero cost.
                            live_items = []
                            for url, v in items:
                                vv = str(v).upper().strip()
                                if vv in cached["active"]:
                                    vdp_results[url] = "Available"
                                elif cached["complete"] or vv in cached["absent"]:
                                    vdp_results[url] = "SOLD (Not in Market Inventory)"
                                else:
                                    live_items.append((url, v))
                                    continue
                                resolved += 1
                                cache_hits += 1
                            if not live_items:
                                continue
                            items = live_items   # only unknowns go to live lookups
                        year_items = [(url, v, get_year(url)) for url, v in items]
                        dom_out = resolve_dealer_marketcheck(dom, year_items, mc_key, gov,
                                                             cache_out=mc_cache_out)
                        for url, status in dom_out.items():
                            vdp_results[url] = status
                            resolved += 1
            except _MCStop as stop:
                stop_reason = stop.reason
            mc_calls_used = gov.count
            leftover = total_unres - resolved
            if mc_cache_out and MC_CACHE_DAYS > 0 and gsheets_conn is not None:
                _mc_cache_save(gsheets_conn, mc_cache_out, max_age_days=MC_CACHE_DAYS)
            if cache_hits:
                st.caption(f"🗂️ {cache_hits} vehicle(s) resolved from the dealer-inventory "
                           f"memory (snapshots under {MC_CACHE_DAYS} days old) — 0 lookups spent.")
            err_detail = getattr(gov, "last_error", None)
            if stop_reason == "limit":
                st.warning(f"The market-inventory service declined further lookups (monthly "
                           f"plan limit reached, or the key needs attention). Resolved "
                           f"{resolved} vehicle(s) before stopping; {leftover} left unresolved."
                           + (f" Last response — {err_detail}" if err_detail else ""))
            elif stop_reason == "budget":
                st.warning(f"Reached this run's market-lookup budget of {MC_BUDGET}. Resolved "
                           f"{resolved} vehicle(s); {leftover} left unresolved. The budget can "
                           f"be raised in settings if needed.")
            elif stop_reason == "monthly":
                st.warning(f"Paused market lookups: the month-wide cap of {MC_MONTH_CAP} "
                           f"(MARKETCHECK_MONTHLY_CAP) is reached. Resolved {resolved} "
                           f"vehicle(s); {leftover} left unresolved this run.")
            elif stop_reason == "errors":
                st.warning(f"The market-inventory lookup stopped after repeated errors. "
                           f"Resolved {resolved} vehicle(s); {leftover} left unresolved."
                           + (f" Last error — {err_detail}" if err_detail else ""))
            elif resolved:
                st.success(f"✅ Resolved {resolved} additional vehicle(s) from live market "
                           f"inventory ({mc_calls_used} lookups used).")
                if err_detail:
                    st.caption(f"⚠️ Some lookups fell back to a slower method — last error: "
                               f"{err_detail}")
            elif err_detail:
                st.error(f"Market-inventory lookups are failing — {leftover} vehicle(s) left "
                         f"unresolved. Last error — {err_detail}")
            else:
                st.info("A few dealers couldn't be matched in the market inventory.")


    df_raw['Sold_Status'] = df_raw['Page Url'].map(vdp_results).fillna('N/A')
    df = df_raw.copy()
    
    df['Is Sold'] = df['Sold_Status'].str.startswith('SOLD')
    df['Vehicle Name'] = df['Page Url'].apply(clean_name_universal)
    df['VIN'] = df['Page Url'].apply(extract_vin)
    df['Type'] = df['Page Url'].apply(extract_type)
    
    vdp_mask = df['Category'] == 'VDP'
    df.loc[vdp_mask, 'Category'] = df.loc[vdp_mask, 'Type'] + ' VDP'
    
    df['Est. Value'] = df.apply(estimate_value, axis=1)
    df['Price Tier'] = df['Est. Value'].apply(get_price_tier)
    
    domain_count = df['Dealer'].nunique()
    eastern = pytz.timezone('US/Eastern')
    report_time = datetime.datetime.now(eastern).strftime('%I:%M %p ET')

    if domain_count > 1:
        report_id = f"Auto Group ({domain_count} Sites) - {report_time}"
    else:
        domain = df['Dealer'].iloc[0] if len(df) > 0 else "Unknown Dealer"
        report_id = f"{domain} ({report_time})"
    
    # --- USAGE TRACKER ---
    try:
        if gsheets_conn is not None:
            dealer_domains = sorted(set(
                urlparse(str(u)).netloc.replace('www.', '').lower() for u in vdp_urls
            ))
            usage_df = gsheets_conn.read(worksheet="UsageStats", ttl=0)
            new_usage = pd.DataFrame([{
                "Timestamp": str(datetime.datetime.now(eastern)),
                "Report": report_id,
                "Dealer URLs": ", ".join(dealer_domains),
                "MarketCheck Calls": mc_calls_used,
            }])
            updated_usage = pd.concat([usage_df, new_usage], ignore_index=True)
            gsheets_conn.update(worksheet="UsageStats", data=updated_usage)
            st.session_state.global_usage_count = len(updated_usage)
            if isinstance(st.session_state.get('mc_month_usage'), int):
                st.session_state.mc_month_usage += mc_calls_used
            elif mc_calls_used:
                st.session_state.mc_month_usage = mc_calls_used
    except Exception:
        pass 
    
    st.session_state.history[report_id] = df
    st.session_state.current_report_id = report_id
    
    st.rerun()

# --- SIDEBAR: SESSION HISTORY MANAGER ---
if st.session_state.history:
    st.sidebar.divider()
    st.sidebar.markdown("### 📂 Session History")
    
    report_names = list(st.session_state.history.keys())
    selected_report = st.sidebar.radio("Select a report to view:", options=report_names, index=report_names.index(st.session_state.current_report_id))
    if selected_report != st.session_state.current_report_id:
        st.session_state.current_report_id = selected_report
        st.rerun()
        
    st.sidebar.divider()
    if st.sidebar.button("🗑️ Clear History"):
        st.session_state.history = {}
        st.session_state.current_report_id = None
        st.rerun()
        
    st.sidebar.divider()
    st.sidebar.markdown(f"📈 **Total Global Scans:** `{st.session_state.global_usage_count}`")
    with st.sidebar.expander("📊 Market Lookup Status", expanded=False):
        st.markdown(f"**Lookup key:** {'✅ detected' if get_marketcheck_key() else '❌ not set'}")
        st.markdown(f"**Per-report budget:** `{MC_BUDGET}` lookups")
        st.markdown(f"**Inventory page cap:** `{MC_PAGE_CAP}` rows")
        st.markdown(f"**Cars per lookup:** `{MC_PAGE_SIZE}`")
        if isinstance(st.session_state.get('mc_month_usage'), int):
            cap_note = f" of `{MC_MONTH_CAP}` cap" if MC_MONTH_CAP else ""
            st.markdown(f"**Lookups this month:** `{st.session_state.mc_month_usage}`{cap_note} (from report log)")
        st.caption("Adjustable in settings via MARKETCHECK_BUDGET, "
                   "MARKETCHECK_PAGE_CAP, and MARKETCHECK_PAGE_SIZE.")

    with st.sidebar.expander("🧪 Market Lookup Diagnostic", expanded=False):
        st.caption("Fires 3 tiny probes (~$0.01 total) to show which filters "
                   "this API key honors on the current endpoint. Use a real "
                   "dealer domain and a VIN currently on their site.")
        diag_domain = st.text_input("Dealer domain", value="hamby.com", key="mc_diag_dom")
        diag_vin = st.text_input("A live VIN from that dealer", value="", key="mc_diag_vin")
        if st.button("Run diagnostic", key="mc_diag_btn"):
            dkey = get_marketcheck_key()
            if not dkey:
                st.error("No MarketCheck key configured.")
            else:
                import requests as _rq
                probes = [("source= (rows=1)", {"source": diag_domain.strip()}),
                          ("source= (rows=50, report-style)", {"source": diag_domain.strip(),
                                                              "rows": "50"}),
                          ("vins=", {"vins": diag_vin.strip().upper()}),
                          ("vin=", {"vin": diag_vin.strip().upper()})]
                # Second probe replicates the report's store call exactly:
                # same params AND the same session construction.
                _diag_sess = _rq.Session()
                try:
                    _retry = Retry(total=3, backoff_factor=1,
                                   status_forcelist=[429, 500, 502, 503, 504])
                    _ad = HTTPAdapter(max_retries=_retry)
                    _diag_sess.mount('https://', _ad)
                except Exception:
                    pass
                for label, extra in probes:
                    if "vin" in label and not diag_vin.strip():
                        st.markdown(f"**{label}** — skipped (no VIN entered)")
                        continue
                    try:
                        _getter = _diag_sess.get if "report-style" in label else _rq.get
                        pr = _getter(MC_ENDPOINT, params={"api_key": dkey, "rows": "1",
                                                          "start": "0", **extra}, timeout=20)
                        if pr.status_code == 200:
                            nf = pr.json().get("num_found", "?")
                            st.markdown(f"**{label}** — HTTP 200, num_found = `{nf}`")
                        else:
                            st.markdown(f"**{label}** — HTTP {pr.status_code}: "
                                        f"`{str(pr.text)[:120]}`")
                    except Exception as e:
                        st.markdown(f"**{label}** — {type(e).__name__}: {str(e)[:120]}")
                    time.sleep(0.25)
                st.caption("How to read this: a filter is honored when num_found is "
                           "small (dealer-sized for source=, ~1-5 for vin=). A number "
                           "in the millions means that filter is being ignored.")

# --- MAIN DASHBOARD DISPLAY ---
if st.session_state.current_report_id is not None:
    st.subheader(f"Viewing Report: {st.session_state.current_report_id}")
    
    df = st.session_state.history[st.session_state.current_report_id].copy()
    
    # 1. IN-APP ACTION CENTER FOR INVENTORY SYNC
    error_df = df[df['Sold_Status'].str.startswith('ERROR', na=False)].copy()
    if not error_df.empty:
        error_df['_base_domain'] = error_df['Page Url'].apply(lambda x: urlparse(str(x)).netloc.replace('www.', '').lower())
        troubleshoot_df = error_df.groupby(['Dealer', '_base_domain']).size().reset_index(name='Blocked Pages')
        severe_blocks = troubleshoot_df[troubleshoot_df['Blocked Pages'] >= 10]
        
        pending_blocks = severe_blocks[~severe_blocks['_base_domain'].isin(DEALER_API_VAULT.keys())]
        
        if len(severe_blocks) > 0:
            if len(pending_blocks) > 0:
                severe_count = len(pending_blocks)
                alert_text = f"⚙️ Action Required: Sync Inventory Database{'s' if severe_count > 1 else ''} ({severe_count} Remaining)"
                
                with st.expander(alert_text, expanded=True):
                    st.markdown("A few dealers hid their inventory behind a firewall. Sync them below to reveal true sales — easiest method first.")

                    st.markdown("---")
                    for _, row in pending_blocks.iterrows():
                        d_name = row['Dealer']
                        d_domain = row['_base_domain']
                        d_blocked = row['Blocked Pages']

                        st.markdown(f"👉 **[{d_name}](https://www.{d_domain})** *({d_blocked} unseen pages)*")
                        ca, cb = st.columns([5, 1])
                        with ca:
                            st.text_input(f"Keys for {d_domain}", key=f"inp_{d_domain}",
                                          label_visibility="collapsed",
                                          placeholder="Paste the line from the 🔑 bookmark (or an Algolia URL)...")
                        with cb:
                            if st.button("💾 Save", key=f"btn_{d_domain}", use_container_width=True):
                                val = st.session_state[f"inp_{d_domain}"]
                                if val:
                                    creds = parse_algolia_credentials(val, fallback_domain=d_domain)
                                    if creds:
                                        dom = creds['domain'] or d_domain
                                        if upsert_vault_row(dom, creds['app_id'], creds['api_key'], creds['index']):
                                            st.rerun()
                                        else:
                                            st.error("Couldn't write to the Google Sheet.")
                                    else:
                                        st.error("Couldn't find all 3 keys in that. Use the 🔑 bookmark on the dealer's site.")
                                else:
                                    st.warning("Paste the keys first!")

                        cc, cd = st.columns(2)
                        with cc:
                            if st.button("⚡ Auto-Detect", key=f"auto_{d_domain}", use_container_width=True,
                                         help="Let the app try to find the keys itself (a strict firewall may block this)."):
                                with st.spinner(f"Searching {d_domain}..."):
                                    creds = autodetect_algolia(d_domain)
                                if creds:
                                    if upsert_vault_row(d_domain, creds['app_id'], creds['api_key'], creds['index']):
                                        st.rerun()
                                    else:
                                        st.error("Found keys but couldn't write to the sheet.")
                                else:
                                    st.warning("Auto-detect blocked. Use the 🔑 bookmark below instead.")
                        with cd:
                            if st.button("🚫 Ignore", key=f"ign_{d_domain}",
                                         help="Mark as not scannable — the dealer doesn't use Dealer Inspire / Algolia, or can't be reached. Stops prompting for it.",
                                         use_container_width=True):
                                if upsert_vault_row(d_domain, 'IGNORE', 'IGNORE', 'IGNORE'):
                                    st.rerun()
                                else:
                                    st.error("Couldn't write the IGNORE flag to the Google Sheet.")
                        st.markdown("")

                    st.markdown("---")
                    st.markdown("**🔑 Easiest way — the one-click key grabber:**")
                    safe_href_ac = html.escape(BOOKMARKLET_JS, quote=True)
                    components.html(
                        f"""
                        <div style="font-family:sans-serif;padding:2px 0;">
                          <a href="{safe_href_ac}"
                             style="display:inline-block;background:#D70015;color:#fff;text-decoration:none;
                                    font-weight:bold;padding:8px 14px;border-radius:6px;border:1px solid #A30010;
                                    cursor:grab;font-size:14px;">🔑 Get Dealer Keys</a>
                          <span style="font-size:12px;color:#777;margin-left:8px;">☝️ one-time: drag me to your bookmarks bar</span>
                        </div>
                        """,
                        height=56,
                    )
                    st.markdown(
                        "1. **Once:** drag the button above to your bookmarks bar.\n"
                        "2. Click a dealer link above to open their site.\n"
                        "3. Click your **Get Dealer Keys** bookmark. If it says the Search Key is missing, "
                        "type any letter in the dealer's inventory search box, then click the bookmark **again**.\n"
                        "4. Paste the copied line into that dealer's box above and hit **Save**."
                    )

                    with st.popover("Advanced: find it manually (DevTools)"):
                        st.markdown("1. Open the dealer's site.\n2. Right-click → **Inspect**.\n3. **Network** tab → **Fetch/XHR** filter → refresh.\n4. Search `inventory`, click the **`query?`** result (contains 'algolia'), right-click → **Copy URL**, and paste it above.")
                        pdf_path = "Dealer Inspire URL Steps.pdf"
                        if os.path.exists(pdf_path):
                            with open(pdf_path, "rb") as pdf_file:
                                pdf_bytes = pdf_file.read()
                            st.download_button(
                                label="📄 Download Visual Guide (PDF)",
                                data=pdf_bytes,
                                file_name="Dealer_Inspire_URL_Steps.pdf",
                                mime="application/pdf"
                            )
            else:
                st.success("✅ All inventory syncs resolved! Click **Run Diagnostic Analysis** in the sidebar to complete the rescan.")

    # 2. REAL-TIME ATTRIBUTION FILTER UI (Sleek Expander Drawer)
    st.markdown("<br>", unsafe_allow_html=True)
    with st.expander("🎯 Interactive VDP Filter", expanded=False):
        st.markdown(
            "<div style='font-size: 14px; color: #555; margin-bottom: 15px;'>"
            "<i>💡 <b>Share of Demand:</b> A vehicle typically needs ~30 total VDP views to sell. "
            "Adjust the slider to set the minimum visitors required to claim marketing influence. "
            "Driving just 5–15 highly qualified visitors represents a dominant share of the demand needed to move a unit.</i></div>", 
            unsafe_allow_html=True
        )
        st.slider("Minimum Visitors to Claim a Sale", min_value=1, max_value=30, step=1, label_visibility="collapsed", key="min_visitors")
    
    min_visitors = st.session_state.min_visitors
    st.markdown("---")
            
    # 3. APPLY FILTER (The "Velocity" Metric Update)
    vdp_df = df[(df['Category'].str.contains('VDP', na=False)) & (df['Attributed Unique Visitors'] >= min_visitors)]
    sold_df = vdp_df[vdp_df['Is Sold']]
    
    new_vdp_all = vdp_df[vdp_df['Type'] == 'New']
    new_sold = sold_df[sold_df['Type'] == 'New']
    new_ltb = (len(new_sold) / len(new_vdp_all) * 100) if len(new_vdp_all) > 0 else 0
    
    used_vdp_all = vdp_df[vdp_df['Type'] == 'Used']
    used_sold = sold_df[sold_df['Type'] == 'Used']
    used_ltb = (len(used_sold) / len(used_vdp_all) * 100) if len(used_vdp_all) > 0 else 0
    
    m_units = len(sold_df)
    m_rev = sold_df['Est. Value'].sum()
    m_pipe = vdp_df['Est. Value'].sum()
    m_ltb = (len(sold_df)/len(vdp_df)*100 if len(vdp_df)>0 else 0)
    
    if not sold_df.empty:
        avg_v = sold_df['Attributed Unique Visitors'].mean()
        missed_threshold = max(avg_v, min_visitors)
        missed_df = df[(df['Sold_Status'] == 'Available') & (df['Category'].str.contains('VDP', na=False)) & (df['Attributed Unique Visitors'] >= missed_threshold)].sort_values('Attributed Unique Visitors', ascending=False).head(10)
    else:
        missed_df = df[(df['Sold_Status'] == 'Available') & (df['Category'].str.contains('VDP', na=False)) & (df['Attributed Unique Visitors'] >= min_visitors)].sort_values('Attributed Unique Visitors', ascending=False).head(10)

    st.markdown("### 📊 Executive Summary")
    m1, m2, m3, m4 = st.columns(4)
    m1.metric("Units Sold (Attributed)", m_units)
    m2.metric("Est. Revenue Sold", f"${m_rev:,.0f}")
    m3.metric("Total Pipeline Value", f"${m_pipe:,.0f}")
    m4.metric(label="Look-to-Book Ratio", value=f"{m_ltb:.1f}%", delta=f"New: {new_ltb:.1f}% | Used: {used_ltb:.1f}%", delta_color="off")

    st.divider()

    domain_count = df['Dealer'].nunique()
    dealer_group_export = None
    
    if domain_count > 1:
        st.markdown(f"### 🏢 Tier 2 / Auto Group Breakdown ({domain_count} Sites)")
        show_tier2 = st.toggle("Display Individual Dealership Insights", value=True)
        
        t2_visitors = df.groupby('Dealer')['Attributed Unique Visitors'].sum().reset_index(name='Total Visitors')
        t2_vdps = vdp_df.groupby('Dealer').agg(
            VDPs_Shopped=('Category', 'count'),
            Pipeline_Value=('Est. Value', 'sum')
        ).reset_index()
        
        t2_sold = sold_df.groupby('Dealer').agg(
            Units_Sold=('Is Sold', 'count'),
            Est_Rev_Sold=('Est. Value', 'sum')
        ).reset_index()
        
        dealer_group = t2_visitors.merge(t2_vdps, on='Dealer', how='left').merge(t2_sold, on='Dealer', how='left').fillna(0)
        dealer_group = dealer_group[dealer_group['VDPs_Shopped'] > 0]
        
        dealer_group['Look-to-Book (%)'] = dealer_group.apply(
            lambda row: round((row['Units_Sold'] / row['VDPs_Shopped'] * 100), 1) if row['VDPs_Shopped'] > 0 else 0.0,
            axis=1
        )
        
        dealer_group = dealer_group.rename(columns={
            'VDPs_Shopped': 'VDPs Shopped',
            'Units_Sold': 'Units Sold',
            'Est_Rev_Sold': 'Est. Rev Sold',
            'Pipeline_Value': 'Pipeline Value'
        })
        
        dealer_group_export = dealer_group.sort_values('Total Visitors', ascending=False)

        if show_tier2:
            cA, cB = st.columns(2)
            with cA:
                fig_traffic = px.bar(dealer_group_export.head(10), x='Dealer', y='Total Visitors', title='Top Dealers by Traffic')
                st.plotly_chart(fig_traffic, use_container_width=True)
            with cB:
                fig_sales = px.bar(dealer_group_export.sort_values('Units Sold', ascending=False).head(10), x='Dealer', y='Units Sold', title='Top Dealers by Attributed Sales')
                st.plotly_chart(fig_sales, use_container_width=True)
                
            display_group = dealer_group_export.copy()
            display_group['Est. Rev Sold'] = display_group['Est. Rev Sold'].apply(lambda x: f"${x:,.0f}")
            display_group['Pipeline Value'] = display_group['Pipeline Value'].apply(lambda x: f"${x:,.0f}")
            
            st.dataframe(display_group, column_config={
                "Look-to-Book (%)": st.column_config.NumberColumn(format="%.1f%%")
            }, use_container_width=True, hide_index=True)
            
        st.divider()
    
    c1, c2, c3 = st.columns(3)
    with c1:
        st.markdown("**Traffic Mix**")
        traffic_data = df.groupby('Category')['Attributed Unique Visitors'].sum().reset_index().sort_values('Attributed Unique Visitors', ascending=False)
        traffic_data = traffic_data.rename(columns={'Attributed Unique Visitors': 'Unique Visits'})
        fig1 = px.bar(traffic_data, x='Category', y='Unique Visits', labels={'Unique Visits': 'Unique Visits'})
        st.plotly_chart(fig1, use_container_width=True)
    with c2:
        st.markdown("**Sales Mix (New vs Used)**")
        if not sold_df.empty:
            type_counts = sold_df['Type'].value_counts().reset_index()
            type_counts.columns = ['Type', 'Count']
            fig2 = px.pie(type_counts, values='Count', names='Type', color='Type', color_discrete_map={'New':'#4F81BD', 'Used':'#C0504D'}, hover_data=['Count'])
            fig2.update_traces(textposition='inside', textinfo='percent+label')
            st.plotly_chart(fig2, use_container_width=True)
    with c3:
        st.markdown("**Sold Value Tiers**")
        if not sold_df.empty:
            tier_counts = sold_df['Price Tier'].value_counts().reset_index()
            tier_counts.columns = ['Price Tier', 'Count']
            fig3 = px.pie(tier_counts, values='Count', names='Price Tier', color_discrete_sequence=px.colors.qualitative.Set2, hover_data=['Count'])
            fig3.update_traces(textposition='inside', textinfo='percent+label')
            st.plotly_chart(fig3, use_container_width=True)

    st.divider()

    t1, t2 = st.columns(2)
    with t1:
        if not sold_df.empty:
            sold_models = sold_df.copy()
            sold_models['Model_Only'] = sold_models['Vehicle Name'].apply(normalize_model)
            model_counts = sold_models['Model_Only'].value_counts().reset_index()
            model_counts.columns = ['Make/Model', 'Units Sold']
            top_models = model_counts[model_counts['Units Sold'] > 1].head(10)
            if not top_models.empty:
                st.subheader("🏆 Top Sold Models (Aggregated > 1 Unit)")
                st.dataframe(top_models, use_container_width=True, hide_index=True)
        else:
            st.info("No sales identified.")

    with t2:
        if not missed_df.empty:
            missed_models = missed_df.copy()
            missed_models['Model_Only'] = missed_models['Vehicle Name'].apply(normalize_model)
            missed_counts = missed_models['Model_Only'].value_counts().reset_index()
            missed_counts.columns = ['Make/Model', 'Missed Count']
            top_missed = missed_counts[missed_counts['Missed Count'] > 1].head(10)
            if not top_missed.empty:
                st.subheader("⚠️ Top Missed Models (Aggregated > 1 Unit)")
                st.dataframe(top_missed, use_container_width=True, hide_index=True)
        else:
            st.info("No missed opportunities identified.")

    st.divider()

    if not sold_df.empty:
        st.subheader("📄 Top Sold Units (Detail)")
        top_sold = sold_df.sort_values('Attributed Unique Visitors', ascending=False).head(10)
        display_sold = top_sold[['Dealer', 'Vehicle Name', 'Type', 'VIN', 'Attributed Unique Visitors', 'Page Url']].reset_index(drop=True)
        display_sold.index += 1
        st.dataframe(display_sold, column_config={"Page Url": st.column_config.LinkColumn("Link", display_text="Open"), "Attributed Unique Visitors": st.column_config.NumberColumn("Visitors")}, use_container_width=True)

    if not missed_df.empty:
        st.write("") 
        st.subheader("👀 Missed Opportunities (Detail)")
        display_missed = missed_df[['Dealer', 'Vehicle Name', 'Type', 'VIN', 'Attributed Unique Visitors', 'Page Url']].sort_values('Attributed Unique Visitors', ascending=False).head(10).reset_index(drop=True)
        display_missed.index += 1
        st.dataframe(display_missed, column_config={"Page Url": st.column_config.LinkColumn("Link", display_text="Open"), "Attributed Unique Visitors": st.column_config.NumberColumn("Visitors")}, use_container_width=True)

    # --- DEALER DEEP DIVES UI ---
    if domain_count > 1 and dealer_group_export is not None:
        st.divider()
        st.markdown("### 🔍 Individual Dealership Deep Dives")
        show_deep_dives = st.toggle("Display Dealership Deep Dives", value=False)
        
        if show_deep_dives:
            expand_all = st.toggle("Expand All Profiles", value=False)
            st.write("")
            
            for dealer_name in dealer_group_export['Dealer']:
                with st.expander(f"🏢 {dealer_name} - Performance Details", expanded=expand_all):
                    d_sold = sold_df[sold_df['Dealer'] == dealer_name]
                    missed_threshold_local = max(avg_v, min_visitors) if not sold_df.empty else min_visitors
                    d_missed = df[(df['Dealer'] == dealer_name) & (df['Sold_Status'] == 'Available') & (df['Category'].str.contains('VDP', na=False)) & (df['Attributed Unique Visitors'] >= missed_threshold_local)]
                    
                    if d_sold.empty and d_missed.empty:
                        st.info("No actionable sales or high-interest missed opportunities meet the current thresholds for this dealer.")
                        continue
                    
                    c_left, c_right = st.columns(2)
                    with c_left:
                        if not d_sold.empty:
                            d_sold_models = d_sold.copy()
                            d_sold_models['Model_Only'] = d_sold_models['Vehicle Name'].apply(normalize_model)
                            d_m_counts = d_sold_models['Model_Only'].value_counts().reset_index()
                            d_m_counts.columns = ['Make/Model', 'Units Sold']
                            st.markdown("**🏆 Top Sold Models**")
                            st.dataframe(d_m_counts.head(5), use_container_width=True, hide_index=True)
                    with c_right:
                        if not d_missed.empty:
                            d_missed_models = d_missed.copy()
                            d_missed_models['Model_Only'] = d_missed_models['Vehicle Name'].apply(normalize_model)
                            d_m_counts = d_missed_models['Model_Only'].value_counts().reset_index()
                            d_m_counts.columns = ['Make/Model', 'Missed Count']
                            st.markdown("**⚠️ Top Missed Models**")
                            st.dataframe(d_m_counts.head(5), use_container_width=True, hide_index=True)
                    
                    if not d_sold.empty:
                        st.markdown("**📄 Top Sold Units (Detail)**")
                        st.dataframe(d_sold[['Vehicle Name', 'Type', 'VIN', 'Attributed Unique Visitors', 'Page Url']].sort_values('Attributed Unique Visitors', ascending=False).head(10), column_config={"Page Url": st.column_config.LinkColumn("Link", display_text="Open")}, use_container_width=True, hide_index=True)
                    
                    if not d_missed.empty:
                        st.markdown("**👀 Missed Opportunities (Detail)**")
                        st.dataframe(d_missed[['Vehicle Name', 'Type', 'VIN', 'Attributed Unique Visitors', 'Page Url']].sort_values('Attributed Unique Visitors', ascending=False).head(10), column_config={"Page Url": st.column_config.LinkColumn("Link", display_text="Open")}, use_container_width=True, hide_index=True)

    st.divider()
    st.markdown("### 📥 Export Reports")
    
    include_missed_in_pdf = st.checkbox("Include 'Missed Opportunities' in Exports (PDF & PowerPoint)", value=True)
    include_dealer_details = False
    if domain_count > 1:
        include_dealer_details = st.checkbox("Include 'Dealer Deep Dives' in Exports (PDF & PowerPoint)", value=False)
        
    st.write("") 

    # Build the dashboard charts as images once, shared by the PDF and PPTX.
    chart_images = build_summary_chart_images(df, sold_df)
    metrics_bundle = {'units_sold': m_units, 'rev_sold': m_rev, 'pipeline': m_pipe, 'ltb': f"{m_ltb:.1f}", 'new_ltb': f"{new_ltb:.1f}", 'used_ltb': f"{used_ltb:.1f}", 'min_visitors': min_visitors}
    export_missed_df = missed_df if not sold_df.empty else pd.DataFrame()

    ex1, ex2, ex3, ex4 = st.columns(4)
    with ex1:
        pdf_data = create_pdf_report(df, sold_df, metrics_bundle, export_missed_df, include_missed_in_pdf, dealer_group_export, include_dealer_details, chart_images=chart_images)
        st.download_button("📥 Download PDF Summary", data=pdf_data, file_name=f"{st.session_state.current_report_id}_Summary.pdf", mime="application/pdf")
    with ex2:
        pptx_data = build_pptx_report(df, sold_df, metrics_bundle, chart_images,
                                      st.session_state.current_report_id,
                                      missed_df=export_missed_df,
                                      dealer_group=dealer_group_export,
                                      include_dealer_details=include_dealer_details,
                                      include_missed=include_missed_in_pdf)
        if pptx_data:
            st.download_button("📥 Download PowerPoint", data=pptx_data,
                               file_name=f"{st.session_state.current_report_id}_Summary.pptx",
                               mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        else:
            st.button("📥 Download PowerPoint", disabled=True,
                      help="PowerPoint export needs the python-pptx package on the host.")
    with ex3:
        st.download_button("📥 Download Sold List (CSV)", sold_df[['Dealer', 'Vehicle Name', 'VIN', 'Page Url', 'Attributed Unique Visitors']].to_csv(index=False), f"{st.session_state.current_report_id}_Sold.csv", "text/csv")
    with ex4:
        st.download_button("📥 Download Full Analysis (CSV)", df.to_csv(index=False), f"{st.session_state.current_report_id}_Full_Analysis.csv", "text/csv")

    st.divider()
    
    with st.expander("ℹ️ Glossary & Guide: How to read this report"):
        st.markdown("""
        ### **Definitions & Insights**
        
        **1. Units Sold**
        The total count of vehicles identified as "Sold" (removed from inventory) *after* receiving attributed traffic from our campaign. This confirms that the audience we drove to the site actively shopped for cars that moved off the lot.
        
        **2. Estimated Value (Rev & Pipeline)**
        A data-driven approximation of the inventory's dollar value. 
        * **New Cars:** Calculated using Base MSRP for the specific model.
        * **Used Cars:** Calculated using the base MSRP depreciated by age.
        * *Note: This is a directional estimate to gauge "Total Pipeline Power," not accounting for specific trim levels, options, or dealer markups.*
        
        **3. Look-to-Book Ratio (Velocity Metric)**
        The efficiency metric of the inventory. It measures the conversion velocity of the cars we drove traffic to. By using the Interactive VDP Filter, this transforms into a Velocity Metric—proving that highly concentrated traffic yields significantly higher conversion rates.
        * *Formula:* `(Sold VDPs) ÷ (Total Active VDPs)`
        
        **4. Tier 2 / Auto Group Breakdown**
        Aggregates performance across multiple dealerships if a multi-domain report is uploaded. 
        * *Clean Reporting:* Dealerships that received general traffic but had **0 VDPs shopped** are automatically filtered out of this table to keep insights focused on high-intent inventory shoppers.
        
        **5. Top Sold Units & Missed Opportunities ("The Watch List")**
        * **Top Sold:** The specific vehicles that received the highest exposure and subsequently sold.
        * **Missed Opportunities:** Active vehicles receiving **above-average traffic** that haven't sold yet. Audit these VDPs immediately for missing photos, "Call for Price" buttons, or pricing outliers!
        
        **6. Inventory Syncs & Troubleshooting**
        If your report shows **0 Sold** and triggers an "Action Required" alert, the dealer's inventory database requires a direct connection. 
        * *Action:* Check the top of your report for the **Action Required** panel. It will guide you to find the API URL. The app will automatically save it to our permanent database (Vault) so that dealer is seamlessly synced moving forward!
        """)

else:
    st.info("👈 Upload a CSV in the sidebar to begin analysis.")
    st.markdown("---")
    st.markdown("### What's New")
    st.markdown("Here's what's new in the Auto-Sales Intelligence Agent:")
    st.markdown("1. **🔓 Blocked dealer sites resolve automatically:** When a dealer's website blocks our scanner, the tool now falls back to live national market inventory to determine which vehicles sold. No more hunting for API URLs — the need for the manual 'Dealer Inspire Fix' and Vault keys should now be minimal to none. Reports just run.")
    st.markdown("2. **📊 New PowerPoint download:** One click builds a client-ready deck styled to fit seamlessly into our Premion attribution reporting — executive summary KPIs, traffic and sales charts, Top Sold, the Missed Opportunities watch list, per-dealer deep dives for auto groups, and a glossary slide.")
    st.markdown("3. **📄 Upgraded PDF report:** A dashboard-style executive summary banner up top, sleeker charts and tables in the Premion palette (and no more cut-off charts), plus a built-in Glossary & Methodology page — ready to forward straight to a client.")
    st.markdown("4. **🚙 Smarter model roll-ups:** Top Sold and Missed Opportunity *models* now aggregate across trims — six Silverado 1500s show as one line instead of six — while the detail tables keep the full trim-level names.")
    st.markdown("5. **🏢 Auto Groups, everywhere:** Multi-dealer reports carry through to both the PDF and the PowerPoint, with an optional store-by-store deep dive for every rooftop in the group.")
